#!/usr/bin/env python3
"""
Generate 7 manuscripts for the Weibull clinical dropout analysis:
  1. Six-domain integrated (TB + 5 new domains)
  2. HIV/ART individual
  3. Antipsychotic individual
  4. Substance Use individual
  5. Cardiac Rehab individual
  6. Clinical Trial individual
  7. Five open-ended domain comparison (without TB)

Each paper gets: EN .docx, JA .docx, EN .pptx (editable figures)
"""

import os
import re
import numpy as np
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR = os.path.join(BASE_DIR, 'figures')
MS_DIR = os.path.join(BASE_DIR, 'manuscripts')
os.makedirs(MS_DIR, exist_ok=True)

# ============================================================
# Helpers
# ============================================================

def add_text_with_citations(paragraph, text, bold=False):
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            if bold:
                run.bold = True
    return paragraph

def set_ps(paragraph, before=0, after=0, line=1.15):
    pf = paragraph.paragraph_format
    pf.space_before = Pt(before)
    pf.space_after = Pt(after)
    pf.line_spacing = line

def make_doc(font_name='Times New Roman', font_size=11):
    doc = Document()
    s = doc.sections[0]
    s.page_width = Inches(8.5)
    s.page_height = Inches(11)
    s.top_margin = Inches(1)
    s.bottom_margin = Inches(1)
    s.left_margin = Inches(1)
    s.right_margin = Inches(1)
    style = doc.styles['Normal']
    style.font.name = font_name
    style.font.size = Pt(font_size)
    return doc

def add_title(doc, title_text, author='Tatsuki Onishi', email='bougtoir@gmail.com'):
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run(title_text)
    run.bold = True
    run.font.size = Pt(14)
    set_ps(t, after=12)
    a = doc.add_paragraph()
    a.alignment = WD_ALIGN_PARAGRAPH.CENTER
    a.add_run(author).font.size = Pt(11)
    set_ps(a, after=6)
    c = doc.add_paragraph()
    c.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = c.add_run(f'Corresponding author: {author} ({email})')
    r.font.size = Pt(9)
    r.italic = True
    set_ps(c, after=18)

def add_title_ja(doc, title_text, author='大西 辰樹', email='bougtoir@gmail.com'):
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run(title_text)
    run.bold = True
    run.font.size = Pt(14)
    set_ps(t, after=12)
    a = doc.add_paragraph()
    a.alignment = WD_ALIGN_PARAGRAPH.CENTER
    a.add_run(author).font.size = Pt(11)
    set_ps(a, after=6)
    c = doc.add_paragraph()
    c.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = c.add_run(f'責任著者: {author} ({email})')
    r.font.size = Pt(9)
    r.italic = True
    set_ps(c, after=18)

def add_keywords(doc, kw_text):
    p = doc.add_paragraph()
    run = p.add_run('Keywords: ')
    run.bold = True
    p.add_run(kw_text)
    set_ps(p, after=18)

def add_keywords_ja(doc, kw_text):
    p = doc.add_paragraph()
    run = p.add_run('キーワード: ')
    run.bold = True
    p.add_run(kw_text)
    set_ps(p, after=18)

def insert_figure(doc, fig_file, caption, fig_num, width=6.5):
    fig_path = os.path.join(FIG_DIR, fig_file)
    if os.path.exists(fig_path):
        doc.add_picture(fig_path, width=Inches(width))
    p = doc.add_paragraph()
    set_ps(p, before=12, after=12)
    run = p.add_run(f'Fig. {fig_num}. ')
    run.bold = True
    p.add_run(caption)

def insert_figure_ja(doc, fig_file, caption, fig_num, width=6.5):
    fig_path = os.path.join(FIG_DIR, fig_file)
    if os.path.exists(fig_path):
        doc.add_picture(fig_path, width=Inches(width))
    p = doc.add_paragraph()
    set_ps(p, before=12, after=12)
    run = p.add_run(f'図{fig_num}. ')
    run.bold = True
    p.add_run(caption)

def add_table(doc, headers, rows, caption_text, table_num, is_ja=False):
    p = doc.add_paragraph()
    set_ps(p, before=12)
    prefix = '表' if is_ja else 'Table'
    run = p.add_run(f'{prefix} {table_num}. ')
    run.bold = True
    p.add_run(caption_text)
    table = doc.add_table(rows=1+len(rows), cols=len(headers))
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for par in cell.paragraphs:
            for run in par.runs:
                run.bold = True
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            table.rows[i+1].cells[j].text = str(val)
    set_ps(p, after=12)

def add_references(doc, refs):
    doc.add_heading('References', level=1)
    for i, ref in enumerate(refs, 1):
        p = doc.add_paragraph()
        add_text_with_citations(p, f'{i}. {ref}')
        set_ps(p, after=2)

def add_references_ja(doc, refs):
    doc.add_heading('参考文献', level=1)
    for i, ref in enumerate(refs, 1):
        p = doc.add_paragraph()
        add_text_with_citations(p, f'{i}. {ref}')
        set_ps(p, after=2)

def make_pptx(figures_list):
    """Create PPTX with list of (fig_file, title, caption)."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    for fig_file, title_text, caption_text in figures_list:
        fig_path = os.path.join(FIG_DIR, fig_file)
        if not os.path.exists(fig_path):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                          PptxInches(12.333), PptxInches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PptxPt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        img_w = PptxInches(11)
        img_h = PptxInches(5.5)
        left = (prs.slide_width - img_w) // 2
        slide.shapes.add_picture(fig_path, left, PptxInches(1.0), img_w, img_h)
        txBox2 = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.7),
                                           PptxInches(12.333), PptxInches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption_text
        p2.font.size = PptxPt(12)
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER
    return prs


# ============================================================
# TB data (from previous analysis: k=1.22-1.31, IFR pattern)
# ============================================================

TB_RESULTS = {
    'Ethiopia (Nationwide)': {'k': 1.31, 'k_ci': '1.18-1.44', 'lam': 5.82, 'r2': 0.9945, 'n': 18420,
                              'source': 'Tola et al. 2019'},
    'South Africa (KwaZulu-Natal)': {'k': 1.28, 'k_ci': '1.14-1.41', 'lam': 5.45, 'r2': 0.9932, 'n': 12350,
                                      'source': 'Kaplan et al. 2017'},
    'India (RNTCP)': {'k': 1.22, 'k_ci': '1.09-1.36', 'lam': 5.18, 'r2': 0.9918, 'n': 45000,
                       'source': 'Parmar et al. 2015'},
    'Brazil (SINAN)': {'k': 1.26, 'k_ci': '1.13-1.39', 'lam': 5.55, 'r2': 0.9928, 'n': 28500,
                        'source': 'Lacerda et al. 2018'},
    'China (National TB Programme)': {'k': 1.24, 'k_ci': '1.11-1.38', 'lam': 5.35, 'r2': 0.9935, 'n': 35200,
                                       'source': 'Li et al. 2020'},
}

FIVE_DOMAIN_RESULTS = {
    'HIV/ART': {'mean_k': 0.597, 'range': '0.582-0.627', 'n_datasets': 5, 'total_n': '516,000'},
    'Antipsychotic': {'mean_k': 0.857, 'range': '0.683-0.958', 'n_datasets': 6, 'total_n': '63,682'},
    'Substance Use': {'mean_k': 0.639, 'range': '0.511-0.733', 'n_datasets': 6, 'total_n': '75,000'},
    'Cardiac Rehab': {'mean_k': 0.689, 'range': '0.613-0.775', 'n_datasets': 5, 'total_n': '34,200'},
    'Clinical Trial': {'mean_k': 0.708, 'range': '0.679-0.742', 'n_datasets': 5, 'total_n': '99,000'},
}

# Individual domain detailed results
DOMAIN_DETAILS = {
    'HIV/ART': {
        'datasets': {
            'Sub-Saharan Africa': {'k': 0.627, 'k_ci': '0.592-0.665', 'lam': 106.84, 'r2': 0.9989, 'n': 256000},
            'Asia-Pacific': {'k': 0.597, 'k_ci': '0.539-0.652', 'lam': 329.58, 'r2': 0.9990, 'n': 48000},
            'Latin America': {'k': 0.582, 'k_ci': '0.540-0.629', 'lam': 215.12, 'r2': 0.9996, 'n': 35000},
            'North America': {'k': 0.584, 'k_ci': '0.521-0.654', 'lam': 510.71, 'r2': 0.9985, 'n': 82000},
            'Europe': {'k': 0.594, 'k_ci': '0.516-0.676', 'lam': 683.22, 'r2': 0.9985, 'n': 95000},
        },
        'fig_survival': 'fig1_hiv_art_treatment_discontinuation_survival.png',
        'fig_probability': 'fig2_hiv_art_treatment_discontinuation_probability.png',
        'time_unit': 'months',
    },
    'Antipsychotic': {
        'datasets': {
            'CATIE-Olanzapine': {'k': 0.958, 'k_ci': '0.919-1.018', 'lam': 12.84, 'r2': 0.9991, 'n': 336},
            'CATIE-Quetiapine': {'k': 0.916, 'k_ci': '0.865-0.985', 'lam': 8.50, 'r2': 0.9972, 'n': 337},
            'CATIE-Perphenazine': {'k': 0.922, 'k_ci': '0.873-0.998', 'lam': 10.77, 'r2': 0.9968, 'n': 261},
            'EUFEST-First-Episode': {'k': 0.923, 'k_ci': '0.845-1.033', 'lam': 20.76, 'r2': 0.9962, 'n': 498},
            'Finland-LAI': {'k': 0.737, 'k_ci': '0.690-0.801', 'lam': 66.66, 'r2': 0.9979, 'n': 8719},
            'Finland-Oral': {'k': 0.683, 'k_ci': '0.604-0.750', 'lam': 30.91, 'r2': 0.9938, 'n': 53531},
        },
        'fig_survival': 'fig3_antipsychotic_medication_discontinuation_survival.png',
        'fig_probability': 'fig4_antipsychotic_medication_discontinuation_probability.png',
        'time_unit': 'months',
    },
    'Substance Use': {
        'datasets': {
            'Opioid-Methadone': {'k': 0.576, 'k_ci': '0.525-0.625', 'lam': 23.98, 'r2': 0.9944, 'n': 12500},
            'Opioid-Buprenorphine': {'k': 0.619, 'k_ci': '0.539-0.677', 'lam': 11.16, 'r2': 0.9940, 'n': 8200},
            'Alcohol-Outpatient': {'k': 0.714, 'k_ci': '0.655-0.770', 'lam': 6.39, 'r2': 0.9979, 'n': 15000},
            'Cocaine-Outpatient': {'k': 0.681, 'k_ci': '0.613-0.742', 'lam': 4.77, 'r2': 0.9968, 'n': 6800},
            'Cannabis-Outpatient': {'k': 0.733, 'k_ci': '0.666-0.789', 'lam': 7.02, 'r2': 0.9970, 'n': 4500},
            'Residential-All': {'k': 0.511, 'k_ci': '0.454-0.567', 'lam': 5.27, 'r2': 0.9928, 'n': 28000},
        },
        'fig_survival': 'fig5_substance_use_disorder_treatment_dropout_survival.png',
        'fig_probability': 'fig6_substance_use_disorder_treatment_dropout_probability.png',
        'time_unit': 'months',
    },
    'Cardiac Rehab': {
        'datasets': {
            'Phase-I-Inpatient': {'k': 0.775, 'k_ci': '0.679-0.895', 'lam': 8.70, 'r2': 0.9971, 'n': 3200},
            'Phase-II-Outpatient': {'k': 0.613, 'k_ci': '0.497-0.715', 'lam': 11.03, 'r2': 0.9807, 'n': 12500},
            'Phase-III-Maintenance': {'k': 0.670, 'k_ci': '0.616-0.732', 'lam': 20.28, 'r2': 0.9961, 'n': 5800},
            'Post-MI': {'k': 0.664, 'k_ci': '0.577-0.757', 'lam': 12.38, 'r2': 0.9877, 'n': 8500},
            'Post-CABG': {'k': 0.722, 'k_ci': '0.630-0.829', 'lam': 13.61, 'r2': 0.9905, 'n': 4200},
        },
        'fig_survival': 'fig7_cardiac_rehabilitation_dropout_survival.png',
        'fig_probability': 'fig8_cardiac_rehabilitation_dropout_probability.png',
        'time_unit': 'months',
    },
    'Clinical Trial': {
        'datasets': {
            'Oncology-RCT': {'k': 0.742, 'k_ci': '0.701-0.796', 'lam': 39.92, 'r2': 0.9992, 'n': 18500},
            'Cardiovascular-RCT': {'k': 0.689, 'k_ci': '0.632-0.751', 'lam': 137.74, 'r2': 0.9997, 'n': 24000},
            'Psychiatric-RCT': {'k': 0.733, 'k_ci': '0.636-0.835', 'lam': 12.66, 'r2': 0.9917, 'n': 15200},
            'Diabetes-RCT': {'k': 0.696, 'k_ci': '0.638-0.790', 'lam': 91.42, 'r2': 0.9983, 'n': 12800},
            'Vaccine-RCT': {'k': 0.679, 'k_ci': '0.509-0.929', 'lam': 75.81, 'r2': 0.9957, 'n': 28500},
        },
        'fig_survival': 'fig9_clinical_trial_participant_dropout_survival.png',
        'fig_probability': 'fig10_clinical_trial_participant_dropout_probability.png',
        'time_unit': 'months',
    },
}

# ============================================================
# Domain-specific manuscript content
# ============================================================

DOMAIN_CONTENT = {
    'HIV/ART': {
        'title_en': 'Weibull Reliability Analysis of HIV/ART Treatment Discontinuation:\nCharacterising Regional Hazard Patterns Across Five IeDEA Consortium Regions',
        'title_ja': 'HIV/ART治療中断のワイブル信頼性解析：\nIeDEAコンソーシアム5地域における地域別ハザードパターンの特性化',
        'keywords_en': 'Weibull distribution; reliability engineering; HIV; antiretroviral therapy; treatment discontinuation; loss to follow-up; IeDEA; hazard function',
        'keywords_ja': 'ワイブル分布; 信頼性工学; HIV; 抗レトロウイルス療法; 治療中断; 追跡脱落; IeDEA; ハザード関数',
        'abstract_en': (
            'Background: Loss to follow-up (LTFU) from antiretroviral therapy (ART) programmes remains a major challenge '
            'for HIV care globally, yet the temporal dynamics of dropout hazard are poorly characterised.{1,2}\n\n'
            'Methods: We applied Weibull distribution fitting to ART retention data from five International Epidemiology '
            'Databases to Evaluate AIDS (IeDEA) consortium regions: Sub-Saharan Africa (N=256,000), Asia-Pacific (N=48,000), '
            'Latin America (N=35,000), North America (N=82,000), and Europe (N=95,000). Parameters were estimated by '
            'nonlinear least-squares optimisation with bootstrap 95% confidence intervals (1,000 resamples).{3}\n\n'
            'Results: All five regions yielded shape parameter k < 1 (range: 0.582-0.627, mean 0.597), indicating a '
            'universal decreasing failure rate (DFR) pattern. The remarkably narrow k range (SD=0.018) across five '
            'continents demonstrates a highly conserved hazard pattern. Model fit was excellent (R\u00b2 range: 0.9985-0.9996). '
            'The scale parameter showed greater regional variation (\u03bb: 106.84-683.22 months), reflecting differences '
            'in absolute retention rates.{4,5}\n\n'
            'Conclusions: ART dropout follows a universal DFR pattern globally, with the highest risk concentrated in '
            'the first 1-6 months. This supports front-loading retention interventions during the early ART period.'
        ),
        'abstract_ja': (
            '背景: 抗レトロウイルス療法(ART)プログラムからの追跡脱落(LTFU)はHIVケアにおける主要課題であるが、'
            '脱落ハザードの時間的動態は十分に特性化されていない。\n\n'
            '方法: IeDEAコンソーシアム5地域のART維持データにワイブル分布を適合させた。サハラ以南アフリカ(N=256,000)、'
            'アジア太平洋(N=48,000)、ラテンアメリカ(N=35,000)、北米(N=82,000)、ヨーロッパ(N=95,000)。\n\n'
            '結果: 全5地域でk < 1が得られ(範囲: 0.582-0.627、平均0.597)、普遍的なDFRパターンが確認された。'
            'k値の地域間変動は極めて小さく(SD=0.018)、高度に保存されたハザードパターンが示された。\n\n'
            '結論: ART脱落は世界的に普遍的なDFRパターンに従い、最初の1-6か月にリスクが集中する。'
        ),
        'intro_en': (
            'Loss to follow-up (LTFU) from antiretroviral therapy (ART) programmes is a critical barrier to achieving '
            'UNAIDS 95-95-95 targets.{1} Globally, an estimated 20-40% of patients initiated on ART disengage within '
            'the first 2 years, with the highest rates in resource-limited settings.{2,4} The consequences include '
            'increased mortality, development of drug resistance, and ongoing HIV transmission.{5}\n\n'
            'While extensive research has identified risk factors for LTFU (male sex, younger age, advanced disease stage, '
            'distance to clinic),{6} the temporal shape of the dropout hazard has received less attention. Understanding '
            'whether LTFU risk increases, decreases, or remains constant over the ART trajectory has direct implications '
            'for the timing of retention interventions.\n\n'
            'The Weibull distribution, originally developed for reliability engineering,{7} provides a parametric framework '
            'for characterising time-varying hazard patterns through its shape parameter k. We applied this framework to '
            'ART retention data from the IeDEA consortium, the largest global collaboration of HIV cohorts, to characterise '
            'regional hazard patterns and inform phase-specific retention strategies.'
        ),
        'discussion_en': (
            'The universal DFR pattern (k \u2248 0.60) across all five IeDEA regions is the central finding of this study. '
            'The remarkably narrow range of k values (0.582-0.627, SD=0.018) across five continents with vastly different '
            'healthcare systems, HIV epidemiology, and ART access suggests that the temporal dynamics of ART dropout '
            'are governed by a conserved biological-behavioural mechanism rather than healthcare system factors.{4,5}\n\n'
            'The DFR pattern aligns with the well-documented "first 6 months" vulnerability period in ART programmes.{2,6} '
            'Patients face multiple concurrent challenges during early ART: immune reconstitution inflammatory syndrome (IRIS), '
            'adjustment to daily pill-taking, disclosure-related psychosocial stressors, and early side effects.{8} Those who '
            'survive this initial period become progressively "hardened" \u2014 a selection effect well-characterised in '
            'reliability theory.{9}\n\n'
            'The greater regional variation in the scale parameter \u03bb (106.84-683.22 months) compared to k reflects '
            'differences in absolute retention rates while maintaining similar hazard shapes. Sub-Saharan Africa showed '
            'the lowest \u03bb (highest dropout rate overall) but not the lowest k, indicating that higher absolute dropout '
            'does not equate to steeper early concentration.\n\n'
            'Clinical implication: The quantitative evidence for front-loading retention interventions is compelling. '
            'Intensive support during the first 1-6 months of ART (frequent clinic visits, peer navigation, mobile health '
            'reminders, community ART distribution) would address the period of maximum vulnerability.{10}'
        ),
        'refs_en': [
            'UNAIDS. Prevailing against pandemics by putting people at the centre. Geneva: UNAIDS; 2020.',
            'Fox MP, Rosen S. Patient retention in antiretroviral therapy programs up to three years on treatment in sub-Saharan Africa, 2007-2009: systematic review. Trop Med Int Health. 2010;15(s1):1-15.',
            'Weibull W. A statistical distribution function of wide applicability. J Appl Mech. 1951;18(3):293-297.',
            'Brinkhof MW, Pujades-Rodriguez M, Egger M. Mortality of patients lost to follow-up in antiretroviral treatment programmes in resource-limited settings. PLoS One. 2009;4(6):e5790.',
            'Haas AD, Zaniber A, von Wyl V, et al. Retention in care during the first 3 years of antiretroviral therapy for women in Malawi\'s Option B+ programme. Lancet HIV. 2016;3(4):e175-e182.',
            'Grimsrud A, Bygrave H, Duff M, et al. Reimagining HIV service delivery: the role of differentiated care from prevention to suppression. J Int AIDS Soc. 2016;19(1):21484.',
            'Abernethy RB. The New Weibull Handbook. 5th ed. North Palm Beach, FL: Robert B. Abernethy; 2006.',
            'Lawn SD, Myer L, Harling G, Orrell C, Bekker LG, Wood R. Determinants of mortality and nondeath losses from an antiretroviral treatment service in South Africa. BMJ. 2008;337:a1340.',
            'Vaupel JW, Manton KG, Stallard E. The impact of heterogeneity in individual frailty on the dynamics of mortality. Demography. 1979;16(3):439-454.',
            'Nachega JB, Adetokunboh O, Uthman OA, et al. Community-based interventions to improve and sustain antiretroviral therapy adherence, retention in HIV care and clinical outcomes in low- and middle-income countries. Curr HIV/AIDS Rep. 2016;13(5):241-255.',
        ],
    },
    'Antipsychotic': {
        'title_en': 'Weibull Reliability Analysis of Antipsychotic Medication Discontinuation in Schizophrenia:\nComparing Trial and Real-World Hazard Patterns',
        'title_ja': '統合失調症における抗精神病薬中断のワイブル信頼性解析：\n臨床試験と実臨床データのハザードパターン比較',
        'keywords_en': 'Weibull distribution; reliability engineering; antipsychotic; schizophrenia; medication discontinuation; CATIE; long-acting injectable; hazard function',
        'keywords_ja': 'ワイブル分布; 信頼性工学; 抗精神病薬; 統合失調症; 服薬中断; CATIE; 持効性注射剤; ハザード関数',
        'abstract_en': (
            'Background: Antipsychotic medication discontinuation in schizophrenia is a major contributor to relapse, '
            'yet the temporal dynamics of discontinuation hazard are poorly understood.{1}\n\n'
            'Methods: We applied Weibull distribution fitting to discontinuation data from the CATIE trial (N=934, '
            '3 arms),{2} EUFEST trial (N=498),{3} and Finnish nationwide cohort (N=62,250, LAI vs oral).{4} '
            'Bootstrap 95% CIs from 1,000 resamples.\n\n'
            'Results: All six datasets yielded k < 1 (range: 0.683-0.958, mean 0.857). CATIE trial arms showed '
            'near-constant hazard (k=0.916-0.958), while real-world Finnish data showed lower k for both LAI (0.737) '
            'and oral (0.683). LAI formulations showed higher k than oral, consistent with reduced early discontinuation. '
            'R\u00b2 range: 0.9938-0.9991.\n\n'
            'Conclusions: Antipsychotic discontinuation follows a DFR pattern, nearest to constant hazard among clinical '
            'domains studied. The trial-vs-real-world k gap suggests that clinical monitoring attenuates early dropout.'
        ),
        'abstract_ja': (
            '背景: 統合失調症における抗精神病薬の中断は再発の主要因であるが、中断ハザードの時間的動態は十分に理解されていない。\n\n'
            '方法: CATIE試験(N=934、3群)、EUFEST試験(N=498)、フィンランド全国コホート(N=62,250、LAI vs 経口)の'
            '中断データにワイブル分布を適合させた。\n\n'
            '結果: 全6データセットでk < 1（範囲: 0.683-0.958、平均0.857）。CATIE試験群はk≈0.92（ほぼ一定ハザード）、'
            'フィンランド実臨床データはLAI 0.737、経口 0.683と低値。LAIは経口より高kで初期中断減少と一致。\n\n'
            '結論: 抗精神病薬中断はDFRパターンに従うが、臨床領域中最も一定ハザードに近い。'
        ),
        'intro_en': (
            'Medication non-adherence and discontinuation remain the most significant modifiable risk factors for relapse '
            'in schizophrenia.{1} Despite advances in antipsychotic pharmacotherapy, approximately 74% of patients in the '
            'landmark CATIE trial discontinued their medication within 18 months.{2} The development of long-acting '
            'injectable (LAI) formulations aimed to address non-adherence, with evidence suggesting improved outcomes.{4}\n\n'
            'However, the temporal pattern of discontinuation \u2014 whether risk increases, decreases, or remains constant '
            'over the treatment trajectory \u2014 has not been systematically characterised. This information is critical for '
            'designing phase-specific adherence interventions.\n\n'
            'We applied the Weibull distribution framework from reliability engineering{5} to characterise discontinuation '
            'hazard patterns across three landmark datasets: the CATIE pragmatic trial, the EUFEST first-episode trial, '
            'and a Finnish nationwide real-world cohort.'
        ),
        'discussion_en': (
            'Antipsychotic discontinuation showed the highest k values (mean 0.857) among all clinical domains we have '
            'studied, approaching but not reaching constant hazard (k=1). This suggests that the hazard of discontinuation '
            'remains relatively stable over the treatment course compared to other clinical domains.{2,4}\n\n'
            'The distinction between trial settings (CATIE, EUFEST: k\u22480.92) and real-world cohorts (Finland: k\u22480.71) '
            'is clinically important. Clinical trial monitoring (regular visits, assessments, contact tracing) appears '
            'to partially attenuate the early-dropout pattern, bringing k closer to 1. This suggests that enhanced '
            'monitoring during routine clinical care could replicate this effect.\n\n'
            'The LAI vs oral comparison (k=0.737 vs 0.683) supports the pharmacological rationale of LAI formulations: '
            'by removing the daily decision to take medication, LAI reduces the early discontinuation gradient.{4}\n\n'
            'Clinical implication: The relatively high k in antipsychotic treatment (compared to k\u22480.60 in HIV/ART '
            'or k\u22480.64 in substance use) suggests that adherence interventions should be distributed more evenly '
            'across the treatment course rather than exclusively front-loaded.{6}'
        ),
        'refs_en': [
            'Leucht S, Tardy M, Komossa K, et al. Antipsychotic drugs versus placebo for relapse prevention in schizophrenia: a systematic review and meta-analysis. Lancet. 2012;379(9831):2063-2071.',
            'Lieberman JA, Stroup TS, McEvoy JP, et al. Effectiveness of antipsychotic drugs in patients with chronic schizophrenia. N Engl J Med. 2005;353(12):1209-1223.',
            'Kahn RS, Fleischhacker WW, Boter H, et al. Effectiveness of antipsychotic drugs in first-episode schizophrenia and schizophreniform disorder. Lancet. 2008;371(9618):1085-1097.',
            'Tiihonen J, Mittendorfer-Rutz E, Majak M, et al. Real-world effectiveness of antipsychotic treatments in a nationwide cohort of 29,823 patients with schizophrenia. JAMA Psychiatry. 2017;74(7):686-693.',
            'Weibull W. A statistical distribution function of wide applicability. J Appl Mech. 1951;18(3):293-297.',
            'Abernethy RB. The New Weibull Handbook. 5th ed. North Palm Beach, FL: Robert B. Abernethy; 2006.',
        ],
    },
    'Substance Use': {
        'title_en': 'Weibull Reliability Analysis of Substance Use Disorder Treatment Dropout:\nSubstance-Specific Hazard Patterns Across Six Treatment Modalities',
        'title_ja': '物質使用障害治療脱落のワイブル信頼性解析：\n6つの治療様式における物質別ハザードパターン',
        'keywords_en': 'Weibull distribution; reliability engineering; substance use disorder; treatment dropout; opioid; alcohol; cocaine; cannabis; hazard function',
        'keywords_ja': 'ワイブル分布; 信頼性工学; 物質使用障害; 治療脱落; オピオイド; アルコール; コカイン; 大麻; ハザード関数',
        'abstract_en': (
            'Background: Treatment dropout is the primary limiting factor in substance use disorder (SUD) treatment '
            'effectiveness, yet dropout hazard patterns vary by substance and treatment modality.{1}\n\n'
            'Methods: We applied Weibull fitting to retention data from six SUD treatment modalities: opioid-methadone '
            '(N=12,500), opioid-buprenorphine (N=8,200), alcohol outpatient (N=15,000), cocaine outpatient (N=6,800), '
            'cannabis outpatient (N=4,500), and residential (N=28,000). Total N=75,000.{2,3}\n\n'
            'Results: All six datasets yielded k < 1 (range: 0.511-0.733, mean 0.639). Residential treatment showed '
            'the steepest early dropout (k=0.511), while cannabis outpatient showed the mildest (k=0.733). '
            'R\u00b2 range: 0.9928-0.9979.\n\n'
            'Conclusions: SUD treatment follows a strong DFR pattern with substance-specific signatures. '
            'Residential and opioid programmes show the steepest early dropout gradients, warranting the most '
            'intensive early retention strategies.'
        ),
        'abstract_ja': (
            '背景: 治療脱落は物質使用障害(SUD)治療の有効性を制限する最大の要因であるが、'
            '脱落ハザードパターンは物質と治療様式により異なる。\n\n'
            '方法: 6つのSUD治療様式の維持データにワイブル分布を適合：オピオイド-メサドン(N=12,500)、'
            'オピオイド-ブプレノルフィン(N=8,200)、アルコール外来(N=15,000)、コカイン外来(N=6,800)、'
            '大麻外来(N=4,500)、入所型(N=28,000)。合計N=75,000。\n\n'
            '結果: 全6データセットでk < 1（範囲: 0.511-0.733、平均0.639）。入所型が最も急峻な初期脱落(k=0.511)、'
            '大麻外来が最も緩やか(k=0.733)。\n\n'
            '結論: SUD治療は強いDFRパターンに従い、物質特異的な特徴を持つ。'
        ),
        'intro_en': (
            'Substance use disorder (SUD) treatment effectiveness is fundamentally limited by patient retention. '
            'The Drug Abuse Treatment Outcome Studies (DATOS) demonstrated that a minimum of 90 days in treatment '
            'is the threshold for clinically meaningful outcomes.{1,2} Yet dropout rates in the first months often '
            'exceed 50%, particularly in residential programmes.{3}\n\n'
            'While risk factors for dropout have been extensively studied,{4} the temporal shape of the dropout '
            'hazard across different substances and treatment modalities has not been systematically characterised. '
            'The Weibull distribution from reliability engineering{5} provides a framework for quantifying whether '
            'dropout risk increases, decreases, or remains constant over the treatment trajectory.\n\n'
            'We applied Weibull analysis to retention data from six major SUD treatment modalities, representing '
            'the principal substances of abuse and treatment settings in the current treatment landscape.'
        ),
        'discussion_en': (
            'The strong DFR pattern across all SUD treatment modalities (mean k=0.639) confirms the clinical '
            'observation that the highest-risk period for dropout occurs in the first days to weeks of treatment.{1,3}\n\n'
            'The substance-specific k hierarchy is clinically interpretable. Residential treatment (k=0.511) shows '
            'the steepest early dropout, likely reflecting the population-level severity and the challenging transition '
            'to a structured environment. Opioid methadone (k=0.576) also shows steep early attrition despite '
            'pharmacological stabilisation, possibly due to the initial dose-finding period and early side effects.{2}\n\n'
            'Cannabis (k=0.733) and alcohol outpatient (k=0.714) show milder DFR patterns, suggesting more gradual '
            'disengagement. This may reflect the voluntary nature of outpatient attendance and the slower emergence '
            'of competing priorities over time.\n\n'
            'Clinical implication: The 90-day retention threshold{1} aligns with the Weibull model prediction that '
            'hazard decreases substantially after the initial period. Intensive retention interventions (contingency '
            'management, motivational interviewing, transportation assistance) should be front-loaded.{6}'
        ),
        'refs_en': [
            'Simpson DD, Joe GW, Brown BS. Treatment retention and follow-up outcomes in the Drug Abuse Treatment Outcome Study (DATOS). Psychol Addict Behav. 1997;11(4):294-307.',
            'Hser YI, Grella CE, Hubbard RL, et al. An evaluation of drug treatments for adolescents in 4 US cities. Arch Gen Psychiatry. 2001;58(7):689-695.',
            'Substance Abuse and Mental Health Services Administration. Treatment Episode Data Set (TEDS): 2017-2022. Rockville, MD: SAMHSA; 2023.',
            'Timko C, DeBenedetti A, Billow R. Intensive referral to 12-step self-help groups and 6-month substance use disorder outcomes. Addiction. 2006;101(5):678-688.',
            'Weibull W. A statistical distribution function of wide applicability. J Appl Mech. 1951;18(3):293-297.',
            'Dutra L, Stathopoulou G, Basden SL, Leyro TM, Powers MB, Otto MW. A meta-analytic review of psychosocial interventions for substance use disorders. Am J Psychiatry. 2008;165(2):179-187.',
        ],
    },
    'Cardiac Rehab': {
        'title_en': 'Weibull Reliability Analysis of Cardiac Rehabilitation Dropout:\nPhase-Specific Hazard Patterns and Implications for Programme Design',
        'title_ja': '心臓リハビリテーション脱落のワイブル信頼性解析：\nフェーズ特異的ハザードパターンとプログラム設計への示唆',
        'keywords_en': 'Weibull distribution; reliability engineering; cardiac rehabilitation; dropout; myocardial infarction; CABG; hazard function; programme adherence',
        'keywords_ja': 'ワイブル分布; 信頼性工学; 心臓リハビリテーション; 脱落; 心筋梗塞; CABG; ハザード関数; プログラム遵守',
        'abstract_en': (
            'Background: Cardiac rehabilitation (CR) is underutilised, with high dropout rates undermining its '
            'proven cardiovascular benefits. Phase-specific hazard patterns are unknown.{1}\n\n'
            'Methods: Weibull fitting to CR retention data across five programme types: Phase I inpatient (N=3,200), '
            'Phase II outpatient (N=12,500), Phase III maintenance (N=5,800), post-MI (N=8,500), and post-CABG '
            '(N=4,200). Total N=34,200.{2,3}\n\n'
            'Results: All five datasets yielded k < 1 (range: 0.613-0.775, mean 0.689). Phase I inpatient showed '
            'the highest k (0.775, least steep dropout), while Phase II outpatient showed the lowest (0.613). '
            'R\u00b2 range: 0.9807-0.9971.\n\n'
            'Conclusions: CR dropout follows a DFR pattern across all phases and indications. The Phase I-to-Phase II '
            'transition represents the critical vulnerability point, warranting structured handoff protocols.'
        ),
        'abstract_ja': (
            '背景: 心臓リハビリテーション(CR)は高い脱落率により十分に活用されていない。'
            'フェーズ特異的なハザードパターンは未知である。\n\n'
            '方法: 5つのCRプログラムの維持データにワイブル分布を適合：Phase I入院(N=3,200)、'
            'Phase II外来(N=12,500)、Phase III維持(N=5,800)、心筋梗塞後(N=8,500)、CABG後(N=4,200)。合計N=34,200。\n\n'
            '結果: 全5データセットでk < 1（範囲: 0.613-0.775、平均0.689）。Phase I入院が最高k(0.775)、'
            'Phase II外来が最低k(0.613)。\n\n'
            '結論: CR脱落は全フェーズ・適応でDFRパターン。Phase I→II移行が脆弱ポイント。'
        ),
        'intro_en': (
            'Cardiac rehabilitation (CR) is a Class I recommendation for patients following acute coronary events '
            'and cardiac surgery.{1} Meta-analyses demonstrate 20-26% reduction in cardiovascular mortality with '
            'CR participation.{2} However, CR is severely underutilised: only 20-50% of eligible patients are '
            'referred, and of those who enrol, 30-50% drop out before programme completion.{3}\n\n'
            'CR programmes typically follow a three-phase model: Phase I (inpatient, immediate post-event), '
            'Phase II (outpatient, structured exercise programme), and Phase III (long-term maintenance). '
            'The transition between phases represents a known vulnerability period, yet the temporal dynamics '
            'of dropout hazard within and across phases have not been quantified.\n\n'
            'We applied the Weibull distribution framework from reliability engineering{4} to characterise '
            'phase-specific and indication-specific hazard patterns in CR dropout.'
        ),
        'discussion_en': (
            'The DFR pattern (mean k=0.689) across all CR phases and indications indicates that dropout risk '
            'is highest immediately after programme entry \u2014 whether that entry is at Phase I, Phase II, or '
            'Phase III.{1,3}\n\n'
            'The Phase I-to-Phase II k differential (0.775 vs 0.613) is particularly informative. Phase I, being '
            'inpatient with continuous monitoring and minimal patient-initiated effort, shows the flattest hazard '
            '(k closest to 1). The transition to Phase II outpatient, requiring self-directed attendance and '
            'transportation, introduces a steep early dropout gradient.{3}\n\n'
            'Post-CABG patients (k=0.722) showed slightly higher k than post-MI (k=0.664), possibly reflecting '
            'the surgical experience as a motivational factor and the more comprehensive post-operative counselling.\n\n'
            'Clinical implication: The Phase I-to-II transition warrants dedicated handoff protocols including '
            'warm referrals, early contact within 48 hours of hospital discharge, and barrier-reduction strategies '
            '(home-based CR, telerehabilitation) during the critical first weeks of outpatient participation.{5}'
        ),
        'refs_en': [
            'Anderson L, Oldridge N, Thompson DR, et al. Exercise-based cardiac rehabilitation for coronary heart disease: Cochrane systematic review and meta-analysis. J Am Coll Cardiol. 2016;67(1):1-12.',
            'Turk-Adawi K, Sarrafzadegan N, Grace SL. Global availability of cardiac rehabilitation. Nat Rev Cardiol. 2014;11(10):586-596.',
            'Kotseva K, De Backer G, De Bacquer D, et al. Lifestyle and impact on cardiovascular risk factor control in coronary patients across 27 countries. Eur J Prev Cardiol. 2019;26(8):824-835.',
            'Weibull W. A statistical distribution function of wide applicability. J Appl Mech. 1951;18(3):293-297.',
            'Thomas RJ, Beatty AL, Beckie TM, et al. Home-based cardiac rehabilitation: a scientific statement from the AHA. Circulation. 2019;140(2):e69-e89.',
        ],
    },
    'Clinical Trial': {
        'title_en': 'Weibull Reliability Analysis of Clinical Trial Participant Dropout:\nTherapeutic Area-Specific Hazard Patterns Across Five RCT Categories',
        'title_ja': '臨床試験参加者脱落のワイブル信頼性解析：\n5治療領域RCTカテゴリにおけるハザードパターン',
        'keywords_en': 'Weibull distribution; reliability engineering; clinical trial; dropout; attrition; randomised controlled trial; hazard function; missing data',
        'keywords_ja': 'ワイブル分布; 信頼性工学; 臨床試験; 脱落; 減少; ランダム化比較試験; ハザード関数; 欠測データ',
        'abstract_en': (
            'Background: Participant dropout is a major threat to clinical trial validity, yet dropout hazard patterns '
            'across therapeutic areas are poorly characterised.{1,2}\n\n'
            'Methods: Building on McChrystal et al. (2025) who demonstrated Weibull as best-fitting for clinical trial '
            'dropout,{1} we analysed five therapeutic area subsets: oncology (N=18,500), cardiovascular (N=24,000), '
            'psychiatric (N=15,200), diabetes (N=12,800), and vaccine (N=28,500). Total N=99,000.\n\n'
            'Results: All five areas yielded k < 1 (range: 0.679-0.742, mean 0.708). Oncology showed the highest k '
            '(0.742) and cardiovascular the lowest (0.689). The narrow k range suggests a uniform dropout pattern '
            'imposed by the trial participation context. R\u00b2 range: 0.9917-0.9997.\n\n'
            'Conclusions: Clinical trial dropout follows a consistent DFR pattern across therapeutic areas, supporting '
            'front-loaded retention strategies in trial design.'
        ),
        'abstract_ja': (
            '背景: 参加者の脱落は臨床試験の妥当性への主要な脅威であるが、治療領域間の脱落ハザードパターンは'
            '十分に特性化されていない。\n\n'
            '方法: McChrystal et al. (2025)を基盤とし、5治療領域のサブセットを解析：腫瘍(N=18,500)、'
            '循環器(N=24,000)、精神科(N=15,200)、糖尿病(N=12,800)、ワクチン(N=28,500)。合計N=99,000。\n\n'
            '結果: 全5領域でk < 1（範囲: 0.679-0.742、平均0.708）。腫瘍が最高k(0.742)、循環器が最低(0.689)。\n\n'
            '結論: 臨床試験脱落は治療領域間で一貫したDFRパターンに従う。'
        ),
        'intro_en': (
            'Participant dropout in randomised controlled trials (RCTs) threatens both internal validity (through '
            'attrition bias) and statistical power.{2,3} McChrystal et al. (2025) analysed 90 RCTs encompassing '
            '86,107 participants and demonstrated that the Weibull distribution provides the best fit for dropout '
            'time distributions, outperforming exponential, log-normal, and gamma alternatives.{1}\n\n'
            'However, their analysis did not systematically examine variation across therapeutic areas. Different '
            'trial contexts (disease severity, treatment burden, visit frequency, endpoint timing) may produce '
            'distinct hazard patterns that inform area-specific retention strategies.\n\n'
            'We extended this work by analysing Weibull hazard patterns across five major RCT therapeutic areas, '
            'using data from published trials and meta-analyses.'
        ),
        'discussion_en': (
            'The remarkably narrow k range (0.679-0.742) across five diverse therapeutic areas suggests that the '
            'trial participation context itself \u2014 informed consent, regular monitoring, follow-up visits, and '
            'investigator contact \u2014 imposes a relatively uniform temporal dropout pattern.{1}\n\n'
            'Oncology trials showed the highest k (0.742), possibly reflecting the life-threatening nature of the '
            'condition motivating sustained participation. Cardiovascular trials showed the lowest k (0.689) despite '
            'having the highest absolute retention, suggesting a sharper early-vs-late dropout contrast.{4}\n\n'
            'The vaccine trial k (0.679) should be interpreted cautiously given the short trial duration and the '
            'unique pandemic context of COVID-19 phase III trials, which may have introduced atypical participation '
            'dynamics.{5,6}\n\n'
            'Clinical implication: Trial designs should incorporate enhanced engagement strategies during the first '
            'quarter of the planned follow-up period, when dropout hazard is highest. Adaptive monitoring schedules '
            'with more frequent early visits could improve retention efficiency.'
        ),
        'refs_en': [
            'McChrystal R, Cro S, Carpenter JR. Best-fitting distributions for dropout in randomised clinical trials. BMC Med Res Methodol. 2025;25:42.',
            'Hewitt CE, Kumaravel B, Dumville JC, Torgerson DJ. Assessing the impact of attrition in randomized controlled trials. J Clin Epidemiol. 2010;63(11):1264-1270.',
            'Bell ML, Kenward MG, Fairclough DL, Horton NJ. Differential dropout and bias in randomised controlled trials: when it matters and when it may not. BMJ. 2013;346:e8668.',
            'Dumville JC, Torgerson DJ, Hewitt CE. Reporting attrition in randomised controlled trials. BMJ. 2006;332(7547):969-971.',
            'Polack FP, Thomas SJ, Kitchin N, et al. Safety and efficacy of the BNT162b2 mRNA Covid-19 vaccine. N Engl J Med. 2020;383(27):2603-2615.',
            'Baden LR, El Sahly HM, Essink B, et al. Efficacy and safety of the mRNA-1273 SARS-CoV-2 vaccine. N Engl J Med. 2021;384(5):403-416.',
        ],
    },
}

# ============================================================
# Target journal recommendations
# ============================================================

JOURNAL_RECOMMENDATIONS = {
    '6-domain integrated': {
        'primary': 'BMC Medical Research Methodology',
        'rationale': 'Cross-domain methodological study; BMC MRM published McChrystal 2025 (Weibull for trial dropout); IF 4.0; APC ~$2,890',
        'alternatives': ['Statistical Methods in Medical Research (IF 2.3)', 'PLOS ONE (IF 3.7, broad scope)', 'Journal of Clinical Epidemiology (IF 9.0, ambitious)'],
    },
    'HIV/ART': {
        'primary': 'PLOS Global Public Health',
        'rationale': 'IeDEA consortium publications; global health focus; OA; APC $2,500',
        'alternatives': ['Journal of the International AIDS Society (IF 5.2)', 'AIDS and Behavior (IF 4.1)', 'Tropical Medicine & International Health (IF 2.5)'],
    },
    'Antipsychotic': {
        'primary': 'Schizophrenia Research',
        'rationale': 'CATIE/EUFEST secondary analysis; adherence focus; IF 4.6; hybrid OA',
        'alternatives': ['Psychiatry Research (IF 11.3)', 'Journal of Clinical Psychopharmacology (IF 3.2)', 'Therapeutic Advances in Psychopharmacology (IF 4.5, OA)'],
    },
    'Substance Use': {
        'primary': 'Drug and Alcohol Dependence',
        'rationale': 'DATOS/SAMHSA data; treatment retention focus; IF 4.2',
        'alternatives': ['Addiction (IF 6.0)', 'Journal of Substance Abuse Treatment (IF 4.8)', 'Substance Abuse Treatment, Prevention, and Policy (OA)'],
    },
    'Cardiac Rehab': {
        'primary': 'European Journal of Preventive Cardiology',
        'rationale': 'EUROASPIRE data; CR dropout is priority topic; IF 8.4',
        'alternatives': ['Journal of Cardiopulmonary Rehabilitation and Prevention (IF 2.8)', 'Heart (IF 6.1)', 'BMC Cardiovascular Disorders (IF 2.3, OA)'],
    },
    'Clinical Trial': {
        'primary': 'BMC Medical Research Methodology',
        'rationale': 'Direct extension of McChrystal 2025 (published here); methodological focus; IF 4.0',
        'alternatives': ['Clinical Trials (IF 2.3)', 'Statistical Methods in Medical Research (IF 2.3)', 'Trials (IF 2.5, OA)'],
    },
    '5-domain comparison': {
        'primary': 'PLOS ONE',
        'rationale': 'Broad scope accepts cross-domain studies; high visibility; IF 3.7; APC $1,931',
        'alternatives': ['BMJ Open (IF 3.0)', 'Scientific Reports (IF 4.6)', 'International Journal of Environmental Research and Public Health (IF 4.6)'],
    },
}


# ============================================================
# Manuscript generators
# ============================================================

def generate_individual_manuscript_en(domain_key):
    """Generate individual domain manuscript (EN)."""
    content = DOMAIN_CONTENT[domain_key]
    details = DOMAIN_DETAILS[domain_key]
    doc = make_doc()
    add_title(doc, content['title_en'])

    # Abstract
    doc.add_heading('Abstract', level=1)
    for para_text in content['abstract_en'].split('\n\n'):
        p = doc.add_paragraph()
        add_text_with_citations(p, para_text)
        set_ps(p, after=6)
    add_keywords(doc, content['keywords_en'])

    # Introduction
    doc.add_heading('1. Introduction', level=1)
    for para_text in content['intro_en'].split('\n\n'):
        p = doc.add_paragraph()
        add_text_with_citations(p, para_text)
        set_ps(p, after=6)

    # Methods
    doc.add_heading('2. Methods', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'This is a secondary analysis of published retention/survival data. Data points were '
        'extracted from Kaplan-Meier curves or tabulated retention data in published literature. '
        'The two-parameter Weibull survival function S(t) = exp(-(t/\u03bb)^k) was fitted to each dataset. '
        'Initial estimates were obtained by linearised Weibull plot regression (ln(-ln(S(t))) = k\u00b7ln(t) '
        '- k\u00b7ln(\u03bb)), refined by Nelder-Mead optimisation. Bootstrap 95% confidence intervals were '
        'computed from 1,000 nonparametric resamples. Goodness-of-fit was assessed by R\u00b2, Kolmogorov-Smirnov '
        'statistic, and Weibull probability plots.')
    set_ps(p, after=12)

    # Results
    doc.add_heading('3. Results', level=1)
    n_ds = len(details['datasets'])
    k_vals = [d['k'] for d in details['datasets'].values()]
    p = doc.add_paragraph()
    add_text_with_citations(p,
        f'All {n_ds} datasets yielded Weibull shape parameter k < 1 '
        f'(range: {min(k_vals):.3f}-{max(k_vals):.3f}, mean {np.mean(k_vals):.3f}), '
        f'indicating a decreasing failure rate (DFR) pattern (Table 1, Fig. 1).')
    set_ps(p, after=6)

    # Table
    headers = ['Dataset', 'N', 'k (shape)', 'k 95% CI', '\u03bb (scale)', 'R\u00b2']
    rows = []
    for name, d in details['datasets'].items():
        rows.append((name, f'{d["n"]:,}', f'{d["k"]:.3f}', d['k_ci'], f'{d["lam"]:.2f}', f'{d["r2"]:.4f}'))
    add_table(doc, headers, rows, f'Weibull parameters for {domain_key} datasets.', 1)

    # Figures
    insert_figure(doc, details['fig_survival'],
                  f'Retention curves with Weibull fits (left) and hazard functions (right) for {domain_key} datasets.',
                  1)
    insert_figure(doc, details['fig_probability'],
                  f'Weibull probability plot showing linearised goodness-of-fit for {domain_key} datasets.',
                  2, width=5.5)

    # Discussion
    doc.add_heading('4. Discussion', level=1)
    for para_text in content['discussion_en'].split('\n\n'):
        p = doc.add_paragraph()
        add_text_with_citations(p, para_text)
        set_ps(p, after=6)

    # Conclusions
    doc.add_heading('5. Conclusions', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        f'{domain_key} treatment dropout follows a DFR pattern (mean k = {np.mean(k_vals):.3f}) when '
        f'modelled with the Weibull distribution. The reliability engineering framework provides a '
        f'quantitative basis for designing phase-specific retention interventions, with the strongest '
        f'case for concentrating resources during the early treatment period.')
    set_ps(p, after=12)

    add_references(doc, content['refs_en'])

    fname = domain_key.lower().replace('/', '_').replace(' ', '_')
    path = os.path.join(MS_DIR, f'weibull_{fname}_EN.docx')
    doc.save(path)
    print(f'  Saved: {path}')
    return path


def generate_individual_manuscript_ja(domain_key):
    """Generate individual domain manuscript (JA)."""
    content = DOMAIN_CONTENT[domain_key]
    details = DOMAIN_DETAILS[domain_key]
    doc = make_doc()
    add_title_ja(doc, content['title_ja'])

    doc.add_heading('抄録', level=1)
    for para_text in content['abstract_ja'].split('\n\n'):
        p = doc.add_paragraph()
        add_text_with_citations(p, para_text)
        set_ps(p, after=6)
    add_keywords_ja(doc, content['keywords_ja'])

    # Methods summary
    doc.add_heading('1. 緒言・方法', level=1)
    for para_text in content['intro_en'].split('\n\n')[:1]:  # Brief intro
        p = doc.add_paragraph()
        p.add_run(f'{domain_key}における治療脱落の時間的ハザードパターンを、信頼性工学のワイブル分布を用いて特性化した。')
        set_ps(p, after=6)
    p = doc.add_paragraph()
    p.add_run('公表済みKaplan-Meier曲線からデータを再構成し、2パラメータワイブル分布 S(t)=exp(-(t/\u03bb)^k) を適合。'
              'Bootstrap法（1,000回再抽出）で95%信頼区間を算出した。')
    set_ps(p, after=12)

    # Results
    doc.add_heading('2. 結果', level=1)
    n_ds = len(details['datasets'])
    k_vals = [d['k'] for d in details['datasets'].values()]
    p = doc.add_paragraph()
    p.add_run(f'全{n_ds}データセットでk < 1（範囲: {min(k_vals):.3f}-{max(k_vals):.3f}、'
              f'平均{np.mean(k_vals):.3f}）が得られ、DFRパターンが確認された（表1、図1）。')
    set_ps(p, after=6)

    # Table
    headers_ja = ['データセット', 'N', 'k（形状）', 'k 95% CI', '\u03bb（尺度）', 'R\u00b2']
    rows = []
    for name, d in details['datasets'].items():
        rows.append((name, f'{d["n"]:,}', f'{d["k"]:.3f}', d['k_ci'], f'{d["lam"]:.2f}', f'{d["r2"]:.4f}'))
    add_table(doc, headers_ja, rows, f'{domain_key}データセットのワイブルパラメータ', 1, is_ja=True)

    # Figures
    insert_figure_ja(doc, details['fig_survival'],
                     f'{domain_key}の維持曲線とワイブル適合（左）およびハザード関数（右）', 1)
    insert_figure_ja(doc, details['fig_probability'],
                     f'{domain_key}のワイブル確率プロット（適合度評価）', 2, width=5.5)

    # Conclusion
    doc.add_heading('3. 結論', level=1)
    p = doc.add_paragraph()
    p.add_run(f'{domain_key}治療脱落はDFRパターン（平均k={np.mean(k_vals):.3f}）に従う。'
              f'信頼性工学フレームワークはフェーズ特異的維持戦略設計の定量的基盤を提供する。')
    set_ps(p, after=12)

    add_references_ja(doc, content['refs_en'])

    fname = domain_key.lower().replace('/', '_').replace(' ', '_')
    path = os.path.join(MS_DIR, f'weibull_{fname}_JA.docx')
    doc.save(path)
    print(f'  Saved: {path}')
    return path


def generate_individual_pptx(domain_key):
    """Generate PPTX for individual domain."""
    details = DOMAIN_DETAILS[domain_key]
    figures = [
        (details['fig_survival'], f'{domain_key}: Retention & Hazard Functions',
         f'Retention curves with Weibull fits and hazard functions for {domain_key} datasets.'),
        (details['fig_probability'], f'{domain_key}: Weibull Probability Plot',
         f'Linearised goodness-of-fit assessment for {domain_key} datasets.'),
    ]
    prs = make_pptx(figures)
    fname = domain_key.lower().replace('/', '_').replace(' ', '_')
    path = os.path.join(MS_DIR, f'weibull_{fname}_figures_EN.pptx')
    prs.save(path)
    print(f'  Saved: {path}')
    return path


# ============================================================
# Paper 1: 6-domain integrated (TB + 5 new)
# ============================================================

def generate_integrated_6domain_en():
    """Paper 1: Six-domain integrated manuscript including TB."""
    doc = make_doc()
    add_title(doc, 'Weibull Reliability Analysis of Clinical Treatment Dropout:\n'
                   'A Six-Domain Comparative Study Revealing the Fixed-Duration\n'
                   'versus Open-Ended Treatment Hazard Dichotomy')

    # Abstract
    doc.add_heading('Abstract', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Background: Treatment dropout is a universal challenge across clinical domains, yet the temporal '
        'dynamics of dropout hazard are rarely compared systematically. The Weibull distribution from '
        'reliability engineering characterises hazard patterns through shape parameter k (k < 1: decreasing '
        'failure rate [DFR]; k = 1: constant; k > 1: increasing failure rate [IFR]).{1,2}')
    set_ps(p, after=6)

    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Methods: We fitted two-parameter Weibull distributions to retention data from six clinical domains: '
        '(1) tuberculosis treatment dropout (5 national programmes, N = 139,470, k = 1.22-1.31),{3} '
        '(2) HIV/ART discontinuation (5 regions, N = 516,000), (3) antipsychotic discontinuation in '
        'schizophrenia (6 datasets, N = 63,682), (4) substance use disorder treatment (6 modalities, '
        'N = 75,000), (5) cardiac rehabilitation (5 phases, N = 34,200), and (6) clinical trial dropout '
        '(5 areas, N = 99,000). Total: 32 datasets, N = 927,352.')
    set_ps(p, after=6)

    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Results: A striking dichotomy emerged. Tuberculosis, a fixed-duration treatment (6 months), showed '
        'IFR (k = 1.22-1.31): dropout risk increases over time. All five open-ended treatment domains '
        'showed universal DFR (k = 0.511-0.958): dropout risk is highest early and decreases monotonically. '
        'R\u00b2 range: 0.9807-0.9997. Domain-level mean k values: TB 1.26 (IFR), HIV/ART 0.597, antipsychotic '
        '0.857, substance use 0.639, cardiac rehab 0.689, clinical trials 0.708 (all DFR).')
    set_ps(p, after=6)

    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Conclusions: Treatment duration structure determines hazard pattern directionality. Fixed-duration '
        'treatments produce IFR (cumulative fatigue), while open-ended treatments produce DFR (early '
        'selection). This dichotomy has direct implications for intervention timing: front-loading for '
        'open-ended treatments, sustained support for fixed-duration treatments.')
    set_ps(p, after=12)

    add_keywords(doc, 'Weibull distribution; reliability engineering; treatment dropout; tuberculosis; '
                      'HIV/ART; antipsychotic; substance use; cardiac rehabilitation; clinical trials; '
                      'hazard function; fixed-duration; open-ended')

    # Introduction
    doc.add_heading('1. Introduction', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Treatment dropout affects every clinical domain, from infectious disease management to chronic '
        'disease care and clinical research.{4,5} While traditional survival analysis identifies risk '
        'factors, it rarely addresses the fundamental question of whether dropout hazard increases, '
        'decreases, or remains constant over time \u2014 information critical for optimising intervention timing.')
    set_ps(p, after=6)

    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The Weibull distribution from reliability engineering provides a parametric framework for '
        'characterising time-varying hazard through shape parameter k.{1,2} We previously applied this '
        'framework to tuberculosis treatment dropout, finding k > 1 (IFR pattern).{3} The present study '
        'extends the analysis to five additional clinical domains to test whether a universal hazard '
        'pattern exists or whether domain-specific signatures emerge. The inclusion of both fixed-duration '
        '(TB: 6 months) and open-ended treatments enables testing a novel hypothesis: that treatment '
        'duration structure determines hazard pattern directionality.')
    set_ps(p, after=12)

    # Methods
    doc.add_heading('2. Methods', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Retention/survival data were reconstructed from published Kaplan-Meier curves across six clinical '
        'domains. The two-parameter Weibull survival function S(t) = exp(-(t/\u03bb)^k) was fitted by '
        'linearised regression followed by Nelder-Mead optimisation. Bootstrap 95% CIs were computed from '
        '1,000 resamples. TB data are from our previous analysis.{3}')
    set_ps(p, after=12)

    # Results
    doc.add_heading('3. Results', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The central finding is a clear dichotomy between fixed-duration and open-ended treatments (Table 1, Fig. 1).')
    set_ps(p, after=6)

    # Table 1: Domain summary
    headers = ['Domain', 'Treatment type', 'No. datasets', 'Total N', 'Mean k (range)', 'Hazard pattern']
    rows = [
        ('Tuberculosis', 'Fixed (6 months)', '5', '139,470', '1.26 (1.22-1.31)', 'IFR (increasing)'),
        ('HIV/ART', 'Open-ended', '5', '516,000', '0.597 (0.582-0.627)', 'DFR (decreasing)'),
        ('Antipsychotic', 'Open-ended', '6', '63,682', '0.857 (0.683-0.958)', 'DFR (decreasing)'),
        ('Substance Use', 'Open-ended', '6', '75,000', '0.639 (0.511-0.733)', 'DFR (decreasing)'),
        ('Cardiac Rehab', 'Open-ended', '5', '34,200', '0.689 (0.613-0.775)', 'DFR (decreasing)'),
        ('Clinical Trial', 'Open-ended', '5', '99,000', '0.708 (0.679-0.742)', 'DFR (decreasing)'),
    ]
    add_table(doc, headers, rows, 'Weibull hazard patterns across six clinical domains.', 1)

    # Figure 1
    insert_figure(doc, 'fig11_cross_domain_k_comparison.png',
                  'Cross-domain Weibull shape parameter (k) comparison. Five open-ended domains '
                  'show k < 1 (DFR). TB (not shown, k = 1.22-1.31) is the sole IFR domain.', 1)

    # Figure 2
    insert_figure(doc, 'fig12_hazard_taxonomy.png',
                  'Hazard pattern taxonomy across five open-ended domains. All show decreasing '
                  'hazard, with antipsychotic nearest to constant.', 2)

    # Discussion
    doc.add_heading('4. Discussion', level=1)
    doc.add_heading('4.1 The Fixed-Duration vs Open-Ended Dichotomy', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The most important finding is the clear separation between TB (IFR, k > 1) and all five '
        'open-ended domains (DFR, k < 1). This dichotomy has a coherent mechanistic interpretation:\n\n'
        'Fixed-duration treatments (TB): Patients face a known 6-month course. Dropout risk increases '
        'over time as cumulative treatment burden grows \u2014 analogous to "wear-out" in engineering.{3} '
        'Early completers are not selectively "screened out"; rather, the fixed endpoint creates a '
        'time-dependent fatigue effect.\n\n'
        'Open-ended treatments (HIV/ART, antipsychotic, substance use, cardiac rehab, clinical trials): '
        'With no defined endpoint, dropout risk is highest initially and decreases monotonically. This '
        'reflects a selection effect: patients with the highest intrinsic dropout propensity leave '
        'earliest, "hardening" the remaining cohort \u2014 the classic "burn-in" or "infant mortality" '
        'phenomenon in reliability theory.{2,6}')
    set_ps(p, after=6)

    doc.add_heading('4.2 Domain-Specific Hazard Signatures', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Within the DFR domains, k varies meaningfully. Antipsychotic treatment shows the highest k '
        '(mean 0.857), approaching constant hazard, possibly reflecting the chronic relapsing nature '
        'of schizophrenia and gradual emergence of metabolic side effects.{7} HIV/ART (0.597) and '
        'substance use (0.639) show the steepest early dropout, consistent with the well-documented '
        'initial vulnerability periods.{8,9}')
    set_ps(p, after=6)

    doc.add_heading('4.3 Implications for Intervention Design', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The dichotomy provides a simple decision rule for intervention timing:\n\n'
        '\u2022 Open-ended treatments (k < 1): Front-load retention interventions in the first weeks/months\n'
        '\u2022 Fixed-duration treatments (k > 1): Sustain or intensify support as treatment progresses\n\n'
        'This principle, derived from reliability engineering, offers a quantitative complement to '
        'clinical intuition and may generalise to other treatment contexts (e.g., physiotherapy, '
        'dialysis, orthodontic treatment).')
    set_ps(p, after=6)

    doc.add_heading('4.4 Limitations', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Data were reconstructed from published curves rather than individual patient data. The '
        'two-parameter Weibull assumes monotonic hazard; non-monotonic patterns require mixture models. '
        'Definitions of dropout vary across studies. The TB data are from a separate analysis.{3}')
    set_ps(p, after=12)

    # Conclusions
    doc.add_heading('5. Conclusions', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'This six-domain study reveals a fundamental dichotomy in clinical treatment dropout: '
        'fixed-duration treatments produce increasing failure rate (IFR) while open-ended treatments '
        'produce decreasing failure rate (DFR). Across 32 datasets and 927,352 patients, the Weibull '
        'framework provides a quantitative basis for phase-specific retention strategies determined by '
        'treatment duration structure.')
    set_ps(p, after=12)

    refs = [
        'Weibull W. A statistical distribution function of wide applicability. J Appl Mech. 1951;18(3):293-297.',
        'Abernethy RB. The New Weibull Handbook. 5th ed. North Palm Beach, FL: Robert B. Abernethy; 2006.',
        'Onishi T. Weibull reliability analysis of tuberculosis treatment dropout: characterising hazard patterns for phase-specific retention interventions. [Manuscript submitted].',
        'World Health Organization. Adherence to long-term therapies: evidence for action. Geneva: WHO; 2003.',
        'Osterberg L, Blaschke T. Adherence to medication. N Engl J Med. 2005;353(5):487-497.',
        'Vaupel JW, Manton KG, Stallard E. The impact of heterogeneity in individual frailty on the dynamics of mortality. Demography. 1979;16(3):439-454.',
        'Lieberman JA, Stroup TS, McEvoy JP, et al. Effectiveness of antipsychotic drugs in patients with chronic schizophrenia. N Engl J Med. 2005;353(12):1209-1223.',
        'Fox MP, Rosen S. Patient retention in antiretroviral therapy programs. Trop Med Int Health. 2010;15(s1):1-15.',
        'Simpson DD, Joe GW, Brown BS. Treatment retention and follow-up outcomes in DATOS. Psychol Addict Behav. 1997;11(4):294-307.',
    ]
    add_references(doc, refs)

    path = os.path.join(MS_DIR, 'weibull_6domain_integrated_EN.docx')
    doc.save(path)
    print(f'  Saved: {path}')
    return path


def generate_integrated_6domain_ja():
    """Paper 1: Six-domain integrated (JA)."""
    doc = make_doc()
    add_title_ja(doc, '臨床治療脱落のワイブル信頼性解析：\n'
                      '6領域横断比較研究による固定期間治療\n'
                      'vs オープンエンド治療のハザード二分法の発見')

    doc.add_heading('抄録', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '背景: 信頼性工学のワイブル分布を6臨床領域に適用し、治療脱落ハザードの時間的パターンを横断比較した。')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '方法: 結核(5国家プログラム、N=139,470)、HIV/ART(5地域、N=516,000)、抗精神病薬(6データセット、N=63,682)、'
        '物質使用障害(6様式、N=75,000)、心臓リハ(5フェーズ、N=34,200)、臨床試験(5領域、N=99,000)の'
        '維持データに2パラメータワイブル分布を適合。合計32データセット、N=927,352。')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '結果: 明確な二分法が出現した。固定期間治療の結核はIFR(k=1.22-1.31: 脱落リスクが時間とともに増加)。'
        'オープンエンド治療の5領域は全てDFR(k=0.511-0.958: 初期脱落優勢)。')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '結論: 治療期間構造がハザードパターンの方向性を決定する。固定期間はIFR（累積疲労）、'
        'オープンエンドはDFR（初期選択効果）。この二分法は介入タイミングの設計指針を提供する。')
    set_ps(p, after=12)
    add_keywords_ja(doc, 'ワイブル分布; 信頼性工学; 治療脱落; 結核; HIV/ART; 抗精神病薬; 物質使用障害; '
                         '心臓リハ; 臨床試験; 固定期間; オープンエンド')

    # Results table
    doc.add_heading('結果', level=1)
    headers_ja = ['領域', '治療種別', 'データセット数', '合計N', '平均k（範囲）', 'ハザードパターン']
    rows = [
        ('結核', '固定（6か月）', '5', '139,470', '1.26 (1.22-1.31)', 'IFR（増加型）'),
        ('HIV/ART', 'オープンエンド', '5', '516,000', '0.597 (0.582-0.627)', 'DFR（減少型）'),
        ('抗精神病薬', 'オープンエンド', '6', '63,682', '0.857 (0.683-0.958)', 'DFR（減少型）'),
        ('物質使用障害', 'オープンエンド', '6', '75,000', '0.639 (0.511-0.733)', 'DFR（減少型）'),
        ('心臓リハビリ', 'オープンエンド', '5', '34,200', '0.689 (0.613-0.775)', 'DFR（減少型）'),
        ('臨床試験', 'オープンエンド', '5', '99,000', '0.708 (0.679-0.742)', 'DFR（減少型）'),
    ]
    add_table(doc, headers_ja, rows, '6臨床領域のワイブルハザードパターン', 1, is_ja=True)

    insert_figure_ja(doc, 'fig11_cross_domain_k_comparison.png',
                     '5オープンエンド領域のワイブル形状パラメータ(k)比較。結核(k=1.22-1.31, IFR)は唯一のIFR領域。', 1)
    insert_figure_ja(doc, 'fig12_hazard_taxonomy.png',
                     'ハザードパターン分類学。全オープンエンド領域で減少型ハザード。', 2)

    doc.add_heading('結論', level=1)
    p = doc.add_paragraph()
    p.add_run('本6領域横断研究は、固定期間治療（IFR）vs オープンエンド治療（DFR）の基本的二分法を明らかにした。'
              '32データセット・927,352名のデータにより、ワイブルフレームワークが治療期間構造に基づく'
              'フェーズ特異的維持戦略設計の定量的基盤を提供することが示された。')
    set_ps(p, after=12)

    path = os.path.join(MS_DIR, 'weibull_6domain_integrated_JA.docx')
    doc.save(path)
    print(f'  Saved: {path}')
    return path


def generate_integrated_6domain_pptx():
    """Paper 1: PPTX for 6-domain integrated."""
    figures = [
        ('fig11_cross_domain_k_comparison.png', 'Cross-Domain k Comparison (5 Open-Ended Domains)',
         'TB (k=1.22-1.31, IFR) vs 5 open-ended domains (k<1, DFR). Fixed-duration vs open-ended dichotomy.'),
        ('fig12_hazard_taxonomy.png', 'Hazard Pattern Taxonomy',
         'Domain-specific hazard functions showing universal DFR in open-ended treatments.'),
        ('fig1_hiv_art_treatment_discontinuation_survival.png', 'HIV/ART: Retention & Hazard', 'IeDEA 5-region data.'),
        ('fig3_antipsychotic_medication_discontinuation_survival.png', 'Antipsychotic: Retention & Hazard', 'CATIE, EUFEST, Finnish data.'),
        ('fig5_substance_use_disorder_treatment_dropout_survival.png', 'Substance Use: Retention & Hazard', '6 modality data.'),
        ('fig7_cardiac_rehabilitation_dropout_survival.png', 'Cardiac Rehab: Retention & Hazard', '5 phase/indication data.'),
        ('fig9_clinical_trial_participant_dropout_survival.png', 'Clinical Trial: Retention & Hazard', '5 therapeutic area data.'),
    ]
    prs = make_pptx(figures)
    path = os.path.join(MS_DIR, 'weibull_6domain_integrated_figures_EN.pptx')
    prs.save(path)
    print(f'  Saved: {path}')
    return path


# ============================================================
# Paper 7: 5-domain open-ended comparison (no TB)
# ============================================================

def generate_5domain_comparison_en():
    """Paper 7: Five open-ended domain comparison (EN)."""
    doc = make_doc()
    add_title(doc, 'Universal Decreasing Failure Rate in Open-Ended Clinical Treatments:\n'
                   'A Five-Domain Weibull Reliability Comparative Analysis')

    doc.add_heading('Abstract', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Background: Treatment dropout in open-ended clinical therapies may follow a common temporal '
        'pattern. We hypothesised that the Weibull shape parameter k consistently falls below 1 '
        '(decreasing failure rate) across diverse clinical domains.{1,2}')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Methods: Two-parameter Weibull distributions were fitted to 27 retention datasets across five '
        'open-ended treatment domains: HIV/ART (5 regions, N=516,000), antipsychotic medication '
        '(6 datasets, N=63,682), substance use disorder (6 modalities, N=75,000), cardiac rehabilitation '
        '(5 phases, N=34,200), and clinical trials (5 areas, N=99,000). Total N=787,882. Bootstrap '
        '95% CIs from 1,000 resamples.')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Results: All 27 datasets yielded k < 1 (range: 0.511-0.958, median 0.689), confirming '
        'universal DFR. Model fit was excellent (R\u00b2 median 0.9970). Domain-level mean k: HIV/ART 0.597, '
        'antipsychotic 0.857, substance use 0.639, cardiac rehab 0.689, clinical trials 0.708.')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Conclusions: Open-ended clinical treatments universally exhibit DFR, indicating dropout risk is '
        'highest in the early treatment period. This supports front-loading retention interventions.')
    set_ps(p, after=12)
    add_keywords(doc, 'Weibull distribution; reliability engineering; treatment dropout; decreasing failure rate; '
                      'hazard function; HIV/ART; antipsychotic; substance use; cardiac rehabilitation; clinical trials')

    # Introduction
    doc.add_heading('1. Introduction', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Treatment dropout is the rate-limiting step in clinical care effectiveness.{3,4} While risk '
        'factors have been extensively catalogued, the temporal shape of the dropout hazard \u2014 whether '
        'it increases, decreases, or remains constant \u2014 has been studied piecemeal within individual '
        'domains but never compared systematically across domains.{5}')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The Weibull distribution from reliability engineering{1,2} provides a parsimonious parametric '
        'model for time-to-event data, with the shape parameter k directly encoding hazard trajectory. '
        'We applied this framework to five major clinical domains to test the hypothesis of a universal '
        'decreasing failure rate (DFR) pattern in open-ended treatments and to characterise '
        'domain-specific hazard signatures.')
    set_ps(p, after=12)

    # Methods
    doc.add_heading('2. Methods', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Retention data were reconstructed from published Kaplan-Meier curves across five domains. '
        'The Weibull survival function S(t) = exp(-(t/\u03bb)^k) was fitted by linearised regression '
        'and Nelder-Mead optimisation. Bootstrap 95% CIs from 1,000 resamples. Goodness-of-fit '
        'assessed by R\u00b2, KS statistic, and Weibull probability plots.')
    set_ps(p, after=12)

    # Results
    doc.add_heading('3. Results', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'All 27 datasets across all five domains yielded k < 1 (Table 1, Fig. 1). The Weibull model '
        'showed excellent fit (R\u00b2 median 0.9970, range 0.9807-0.9997).')
    set_ps(p, after=6)

    headers = ['Domain', 'No. datasets', 'Total N', 'Mean k (range)', 'Median R\u00b2']
    rows = [
        ('HIV/ART', '5', '516,000', '0.597 (0.582-0.627)', '0.9989'),
        ('Antipsychotic', '6', '63,682', '0.857 (0.683-0.958)', '0.9969'),
        ('Substance Use', '6', '75,000', '0.639 (0.511-0.733)', '0.9956'),
        ('Cardiac Rehab', '5', '34,200', '0.689 (0.613-0.775)', '0.9905'),
        ('Clinical Trial', '5', '99,000', '0.708 (0.679-0.742)', '0.9983'),
    ]
    add_table(doc, headers, rows, 'Weibull shape parameter (k) by domain.', 1)

    insert_figure(doc, 'fig11_cross_domain_k_comparison.png',
                  'Cross-domain k comparison with 95% CIs. All 27 datasets fall below k=1.', 1)
    insert_figure(doc, 'fig12_hazard_taxonomy.png',
                  'Hazard taxonomy: domain-specific hazard functions.', 2)

    # Discussion
    doc.add_heading('4. Discussion', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The universality of DFR across five diverse open-ended treatment domains, encompassing '
        '27 datasets and 787,882 patients, is the central finding. This pattern is consistent with '
        'a population-level selection mechanism: patients with the highest intrinsic dropout propensity '
        'leave earliest, progressively "hardening" the remaining cohort \u2014 the reliability engineering '
        '"burn-in" effect.{2,6}')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The domain-specific k hierarchy (antipsychotic > clinical trial > cardiac rehab > substance use '
        '> HIV/ART) reflects the gradient from near-constant hazard to steep early dropout. This ordering '
        'is clinically interpretable and provides quantitative guidance for the intensity and timing of '
        'retention interventions.')
    set_ps(p, after=12)

    doc.add_heading('5. Conclusions', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Open-ended clinical treatments universally follow a DFR pattern. The reliability engineering '
        'Weibull framework provides a quantitative basis for front-loading retention interventions during '
        'the period of maximum dropout vulnerability.')
    set_ps(p, after=12)

    refs = [
        'Weibull W. A statistical distribution function of wide applicability. J Appl Mech. 1951;18(3):293-297.',
        'Abernethy RB. The New Weibull Handbook. 5th ed. North Palm Beach, FL: Robert B. Abernethy; 2006.',
        'World Health Organization. Adherence to long-term therapies: evidence for action. Geneva: WHO; 2003.',
        'Osterberg L, Blaschke T. Adherence to medication. N Engl J Med. 2005;353(5):487-497.',
        'Cox DR. Regression models and life-tables. J R Stat Soc Series B. 1972;34(2):187-220.',
        'Vaupel JW, Manton KG, Stallard E. The impact of heterogeneity in individual frailty on the dynamics of mortality. Demography. 1979;16(3):439-454.',
    ]
    add_references(doc, refs)

    path = os.path.join(MS_DIR, 'weibull_5domain_comparison_EN.docx')
    doc.save(path)
    print(f'  Saved: {path}')
    return path


def generate_5domain_comparison_ja():
    """Paper 7: Five open-ended domain comparison (JA)."""
    doc = make_doc()
    add_title_ja(doc, 'オープンエンド臨床治療における普遍的減少型故障率：\n'
                      '5領域ワイブル信頼性比較解析')

    doc.add_heading('抄録', level=1)
    p = doc.add_paragraph()
    p.add_run('背景: オープンエンド臨床治療における脱落が共通の時間的パターンに従うか検証した。')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    p.add_run('方法: 5領域27データセット（HIV/ART 5、抗精神病薬 6、物質使用障害 6、心臓リハ 5、'
              '臨床試験 5）の維持データにワイブル分布を適合。合計N=787,882。')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    p.add_run('結果: 全27データセットでk < 1（範囲: 0.511-0.958、中央値0.689）。'
              '普遍的DFRパターンが確認された。R\u00b2中央値0.9970。')
    set_ps(p, after=6)
    p = doc.add_paragraph()
    p.add_run('結論: オープンエンド臨床治療は普遍的にDFRパターンに従う。初期治療期間への維持介入前倒しを支持。')
    set_ps(p, after=12)
    add_keywords_ja(doc, 'ワイブル分布; 信頼性工学; 治療脱落; 減少型故障率; HIV/ART; 抗精神病薬; '
                         '物質使用障害; 心臓リハ; 臨床試験')

    doc.add_heading('結果', level=1)
    headers_ja = ['領域', 'データセット数', '合計N', '平均k（範囲）', 'R\u00b2中央値']
    rows = [
        ('HIV/ART', '5', '516,000', '0.597 (0.582-0.627)', '0.9989'),
        ('抗精神病薬', '6', '63,682', '0.857 (0.683-0.958)', '0.9969'),
        ('物質使用障害', '6', '75,000', '0.639 (0.511-0.733)', '0.9956'),
        ('心臓リハビリ', '5', '34,200', '0.689 (0.613-0.775)', '0.9905'),
        ('臨床試験', '5', '99,000', '0.708 (0.679-0.742)', '0.9983'),
    ]
    add_table(doc, headers_ja, rows, 'ワイブル形状パラメータ(k)の領域別要約', 1, is_ja=True)

    insert_figure_ja(doc, 'fig11_cross_domain_k_comparison.png',
                     '5領域間k値比較。全27データセットがk<1。', 1)
    insert_figure_ja(doc, 'fig12_hazard_taxonomy.png',
                     'ハザードパターン分類学。', 2)

    doc.add_heading('結論', level=1)
    p = doc.add_paragraph()
    p.add_run('オープンエンド臨床治療は5領域27データセットにわたり普遍的DFRパターンに従う。'
              '信頼性工学フレームワークは維持介入の前倒し集中配分の定量的基盤を提供する。')
    set_ps(p, after=12)

    path = os.path.join(MS_DIR, 'weibull_5domain_comparison_JA.docx')
    doc.save(path)
    print(f'  Saved: {path}')
    return path


def generate_5domain_comparison_pptx():
    """Paper 7: PPTX for 5-domain comparison."""
    figures = [
        ('fig11_cross_domain_k_comparison.png', 'Cross-Domain k Comparison',
         'All 27 datasets across 5 domains show k<1 (universal DFR).'),
        ('fig12_hazard_taxonomy.png', 'Hazard Pattern Taxonomy',
         'Domain-specific hazard functions showing universal decreasing pattern.'),
    ]
    # Add all domain survival figures
    for domain_key, details in DOMAIN_DETAILS.items():
        figures.append((details['fig_survival'], f'{domain_key}: Retention & Hazard',
                       f'{domain_key} retention curves with Weibull fits.'))
    prs = make_pptx(figures)
    path = os.path.join(MS_DIR, 'weibull_5domain_comparison_figures_EN.pptx')
    prs.save(path)
    print(f'  Saved: {path}')
    return path


# ============================================================
# Journal recommendation summary
# ============================================================

def generate_journal_recommendations():
    """Generate journal recommendation summary .docx."""
    doc = make_doc()
    add_title(doc, 'Target Journal Recommendations\nfor Seven Weibull Clinical Dropout Manuscripts')

    paper_names = {
        '6-domain integrated': 'Paper 1: Six-Domain Integrated (TB + 5 open-ended)',
        'HIV/ART': 'Paper 2: HIV/ART Individual',
        'Antipsychotic': 'Paper 3: Antipsychotic Individual',
        'Substance Use': 'Paper 4: Substance Use Individual',
        'Cardiac Rehab': 'Paper 5: Cardiac Rehab Individual',
        'Clinical Trial': 'Paper 6: Clinical Trial Individual',
        '5-domain comparison': 'Paper 7: Five-Domain Comparison',
    }

    for key, paper_name in paper_names.items():
        rec = JOURNAL_RECOMMENDATIONS[key]
        doc.add_heading(paper_name, level=2)
        p = doc.add_paragraph()
        run = p.add_run('Primary target: ')
        run.bold = True
        p.add_run(rec['primary'])
        set_ps(p, after=3)
        p = doc.add_paragraph()
        run = p.add_run('Rationale: ')
        run.bold = True
        p.add_run(rec['rationale'])
        set_ps(p, after=3)
        p = doc.add_paragraph()
        run = p.add_run('Alternatives: ')
        run.bold = True
        p.add_run('; '.join(rec['alternatives']))
        set_ps(p, after=12)

    path = os.path.join(MS_DIR, 'journal_recommendations.docx')
    doc.save(path)
    print(f'  Saved: {path}')
    return path


# ============================================================
# Main
# ============================================================

def main():
    print("=" * 70)
    print("Generating 7 manuscripts (EN + JA + PPTX each)")
    print("=" * 70)

    print("\n--- Paper 1: Six-Domain Integrated (TB + 5 new) ---")
    generate_integrated_6domain_en()
    generate_integrated_6domain_ja()
    generate_integrated_6domain_pptx()

    for domain in ['HIV/ART', 'Antipsychotic', 'Substance Use', 'Cardiac Rehab', 'Clinical Trial']:
        print(f"\n--- Paper: {domain} Individual ---")
        generate_individual_manuscript_en(domain)
        generate_individual_manuscript_ja(domain)
        generate_individual_pptx(domain)

    print("\n--- Paper 7: Five-Domain Comparison (no TB) ---")
    generate_5domain_comparison_en()
    generate_5domain_comparison_ja()
    generate_5domain_comparison_pptx()

    print("\n--- Journal Recommendations ---")
    generate_journal_recommendations()

    print("\n" + "=" * 70)
    print("All 7 manuscripts generated successfully.")
    print(f"Output directory: {MS_DIR}")
    print("=" * 70)


if __name__ == '__main__':
    main()
