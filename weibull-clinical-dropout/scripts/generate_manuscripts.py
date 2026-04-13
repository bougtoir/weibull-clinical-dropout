#!/usr/bin/env python3
"""
Generate manuscripts (.docx) and figure slides (.pptx) for all 5 domains.
Each domain gets:
  - English manuscript (.docx) with inline figures/tables
  - Japanese manuscript (.docx) with inline figures/tables
  - English editable figures (.pptx, 1 figure per slide, widescreen)

Also generates a comprehensive cross-domain comparison manuscript.
"""

import os
import re
import numpy as np
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu as PptxEmu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR = os.path.join(BASE_DIR, 'figures')
MS_DIR = os.path.join(BASE_DIR, 'manuscripts')
os.makedirs(MS_DIR, exist_ok=True)


# ============================================================
# Helper: superscript citations using font-based approach
# ============================================================

def add_text_with_citations(paragraph, text, bold=False):
    """
    Parse text with {ref} markers and create runs with proper superscript.
    E.g., "text.{1-3} more" -> "text." + superscript("1-3") + " more"
    """
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            if bold:
                run.bold = True
    return paragraph


def set_paragraph_spacing(paragraph, space_before=0, space_after=0, line_spacing=1.15):
    """Set paragraph spacing."""
    pf = paragraph.paragraph_format
    pf.space_before = Pt(space_before)
    pf.space_after = Pt(space_after)
    pf.line_spacing = line_spacing


# ============================================================
# Results data (from analysis output)
# ============================================================

RESULTS = {
    'HIV/ART': {
        'Sub-Saharan Africa': {'k': 0.627, 'k_ci': '0.590-0.666', 'lam': 106.84, 'r2': 0.9989, 'n': 256000},
        'Asia-Pacific': {'k': 0.597, 'k_ci': '0.548-0.647', 'lam': 329.58, 'r2': 0.9990, 'n': 48000},
        'Latin America': {'k': 0.582, 'k_ci': '0.542-0.623', 'lam': 215.12, 'r2': 0.9996, 'n': 35000},
        'North America': {'k': 0.584, 'k_ci': '0.517-0.649', 'lam': 510.71, 'r2': 0.9985, 'n': 82000},
        'Europe': {'k': 0.594, 'k_ci': '0.521-0.678', 'lam': 683.22, 'r2': 0.9985, 'n': 95000},
    },
    'Antipsychotic': {
        'CATIE-Olanzapine': {'k': 0.958, 'k_ci': '0.921-1.009', 'lam': 12.84, 'r2': 0.9991, 'n': 336},
        'CATIE-Quetiapine': {'k': 0.916, 'k_ci': '0.860-0.985', 'lam': 8.50, 'r2': 0.9972, 'n': 337},
        'CATIE-Perphenazine': {'k': 0.922, 'k_ci': '0.871-0.993', 'lam': 10.77, 'r2': 0.9968, 'n': 261},
        'EUFEST-First-Episode': {'k': 0.923, 'k_ci': '0.849-1.054', 'lam': 20.76, 'r2': 0.9962, 'n': 498},
        'Finland-LAI': {'k': 0.737, 'k_ci': '0.692-0.796', 'lam': 66.66, 'r2': 0.9979, 'n': 8719},
        'Finland-Oral': {'k': 0.683, 'k_ci': '0.602-0.749', 'lam': 30.91, 'r2': 0.9938, 'n': 53531},
    },
    'Substance Use': {
        'Opioid-Methadone': {'k': 0.576, 'k_ci': '0.522-0.628', 'lam': 23.98, 'r2': 0.9944, 'n': 12500},
        'Opioid-Buprenorphine': {'k': 0.619, 'k_ci': '0.558-0.693', 'lam': 11.16, 'r2': 0.9940, 'n': 8200},
        'Alcohol-Outpatient': {'k': 0.714, 'k_ci': '0.656-0.769', 'lam': 6.39, 'r2': 0.9979, 'n': 15000},
        'Cocaine-Outpatient': {'k': 0.681, 'k_ci': '0.613-0.741', 'lam': 4.77, 'r2': 0.9968, 'n': 6800},
        'Cannabis-Outpatient': {'k': 0.733, 'k_ci': '0.666-0.787', 'lam': 7.02, 'r2': 0.9970, 'n': 4500},
        'Residential-All': {'k': 0.511, 'k_ci': '0.461-0.569', 'lam': 5.27, 'r2': 0.9928, 'n': 28000},
    },
    'Cardiac Rehab': {
        'Phase-I-Inpatient': {'k': 0.775, 'k_ci': '0.679-0.875', 'lam': 8.70, 'r2': 0.9971, 'n': 3200},
        'Phase-II-Outpatient': {'k': 0.613, 'k_ci': '0.503-0.728', 'lam': 11.03, 'r2': 0.9807, 'n': 12500},
        'Phase-III-Maintenance': {'k': 0.670, 'k_ci': '0.618-0.728', 'lam': 20.28, 'r2': 0.9961, 'n': 5800},
        'Post-MI': {'k': 0.664, 'k_ci': '0.582-0.752', 'lam': 12.38, 'r2': 0.9877, 'n': 8500},
        'Post-CABG': {'k': 0.722, 'k_ci': '0.636-0.823', 'lam': 13.61, 'r2': 0.9905, 'n': 4200},
    },
    'Clinical Trial': {
        'Oncology-RCT': {'k': 0.742, 'k_ci': '0.702-0.796', 'lam': 39.92, 'r2': 0.9992, 'n': 18500},
        'Cardiovascular-RCT': {'k': 0.689, 'k_ci': '0.638-0.758', 'lam': 137.74, 'r2': 0.9997, 'n': 24000},
        'Psychiatric-RCT': {'k': 0.733, 'k_ci': '0.656-0.826', 'lam': 12.66, 'r2': 0.9917, 'n': 15200},
        'Diabetes-RCT': {'k': 0.696, 'k_ci': '0.645-0.793', 'lam': 91.42, 'r2': 0.9983, 'n': 12800},
        'Vaccine-RCT': {'k': 0.679, 'k_ci': '0.507-0.939', 'lam': 75.81, 'r2': 0.9957, 'n': 28500},
    },
}


# ============================================================
# Cross-domain comprehensive manuscript (English)
# ============================================================

def generate_comprehensive_manuscript_en():
    """Generate the main cross-domain comparison manuscript in English."""
    doc = Document()
    
    # Page setup
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(11)
    
    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('Weibull Reliability Analysis of Clinical Treatment Dropout:\n'
                        'A Five-Domain Comparative Study Characterising Hazard Patterns\n'
                        'for Phase-Specific Retention Interventions')
    run.bold = True
    run.font.size = Pt(14)
    set_paragraph_spacing(title, space_after=12)
    
    # Authors
    authors = doc.add_paragraph()
    authors.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = authors.add_run('Tatsuki Onishi')
    run.font.size = Pt(11)
    set_paragraph_spacing(authors, space_after=6)
    
    # Corresponding author
    corr = doc.add_paragraph()
    corr.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = corr.add_run('Corresponding author: Tatsuki Onishi (bougtoir@gmail.com)')
    run.font.size = Pt(9)
    run.italic = True
    set_paragraph_spacing(corr, space_after=18)
    
    # ---- ABSTRACT ----
    h = doc.add_heading('Abstract', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Background: Treatment dropout remains a critical challenge across clinical domains, '
        'yet the temporal dynamics of dropout hazard are rarely quantified. Reliability engineering\'s '
        'Weibull distribution provides a parametric framework for characterising time-varying hazard '
        'patterns through the shape parameter k, where k < 1 indicates decreasing failure rate (DFR; '
        'early dropout dominant), k = 1 indicates constant hazard, and k > 1 indicates increasing '
        'failure rate (IFR; wear-out pattern).')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Methods: We applied two-parameter Weibull distribution fitting to retention/survival data '
        'reconstructed from published Kaplan-Meier curves across five clinical domains: (1) HIV/ART '
        'treatment discontinuation (5 regions, N = 516,000), (2) antipsychotic medication '
        'discontinuation in schizophrenia (6 datasets, N = 63,682), (3) substance use disorder '
        'treatment dropout (6 substance types, N = 75,000), (4) cardiac rehabilitation dropout '
        '(5 programme phases, N = 34,200), and (5) clinical trial participant withdrawal '
        '(5 therapeutic areas, N = 99,000). Parameters were estimated by nonlinear least-squares '
        'optimisation with bootstrap 95% confidence intervals (1,000 resamples).')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Results: All 27 datasets across all five domains yielded k < 1 (range: 0.511-0.958), '
        'demonstrating a universal decreasing failure rate pattern. The Weibull model showed '
        'excellent fit (median R\u00b2 = 0.9970, range: 0.9807-0.9997). Domain-level mean k values '
        'were: HIV/ART 0.597 (0.582-0.627), antipsychotic 0.857 (0.683-0.958), substance use '
        '0.639 (0.511-0.733), cardiac rehabilitation 0.689 (0.613-0.775), and clinical trials '
        '0.708 (0.679-0.742). Antipsychotics showed the highest k values (nearest to constant '
        'hazard), while HIV/ART and substance use showed the lowest (steepest early dropout).')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Conclusions: The universal DFR pattern across five disparate clinical domains indicates '
        'that dropout risk is consistently highest in the early treatment period and decreases '
        'monotonically over time. This finding has direct implications for resource allocation: '
        'front-loading retention interventions during the initial treatment phase would address '
        'the period of maximum vulnerability. The reliability engineering framework provides a '
        'quantitative basis for designing phase-specific retention strategies.')
    set_paragraph_spacing(p, space_after=12)
    
    # Keywords
    p = doc.add_paragraph()
    run = p.add_run('Keywords: ')
    run.bold = True
    p.add_run('Weibull distribution; reliability engineering; treatment dropout; loss to follow-up; '
              'hazard function; survival analysis; HIV/ART; antipsychotic; substance use disorder; '
              'cardiac rehabilitation; clinical trials')
    set_paragraph_spacing(p, space_after=18)
    
    # ---- INTRODUCTION ----
    doc.add_heading('1. Introduction', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Treatment dropout, variably termed loss to follow-up (LTFU), treatment discontinuation, '
        'or non-adherence, represents one of the most pervasive challenges in clinical medicine.{1,2} '
        'Across therapeutic domains ranging from infectious disease management to psychiatric care, '
        'substantial proportions of patients disengage from prescribed treatments before achieving '
        'therapeutic goals. The consequences are multifaceted: increased morbidity and mortality, '
        'development of drug resistance, higher healthcare costs, and compromised public health '
        'outcomes.{3,4}')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Traditional survival analysis approaches, including Kaplan-Meier estimation and Cox '
        'proportional hazards models, have been widely applied to characterise treatment dropout.{5} '
        'However, these methods focus primarily on identifying risk factors (covariates) rather '
        'than characterising the temporal shape of the dropout hazard itself. The question of '
        'whether dropout risk increases, decreases, or remains constant over the treatment '
        'trajectory has received insufficient attention, despite its direct relevance to the '
        'timing and intensity of retention interventions.')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Reliability engineering offers a well-established parametric framework for addressing '
        'this gap. The Weibull distribution, developed by Waloddi Weibull in 1951,{6} has been '
        'the cornerstone of failure-time analysis in engineering for over seven decades. Its '
        'two-parameter form, characterised by shape parameter k and scale parameter \u03bb, provides '
        'a parsimonious yet flexible model for time-to-event data. The shape parameter k directly '
        'encodes the hazard pattern: k < 1 produces a decreasing failure rate (DFR), analogous to '
        '"infant mortality" in engineering; k = 1 yields a constant hazard (exponential distribution); '
        'and k > 1 produces an increasing failure rate (IFR), corresponding to "wear-out" failures.{7}')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Recent applications have demonstrated the utility of this framework beyond traditional '
        'engineering contexts. We previously applied Weibull analysis to the reign durations of '
        'historical monarchs,{8} and to tuberculosis treatment dropout,{9} revealing distinct '
        'hazard patterns with direct interpretive value. The present study extends this approach '
        'to a comprehensive cross-domain comparison, applying Weibull analysis to treatment dropout '
        'data from five major clinical areas: HIV/antiretroviral therapy (ART), antipsychotic '
        'medication in schizophrenia, substance use disorder treatment, cardiac rehabilitation, '
        'and clinical trial participation.')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The objectives of this study are: (1) to fit Weibull distributions to retention data '
        'across five clinical domains, (2) to compare hazard patterns (k values) across domains '
        'and sub-populations, (3) to assess whether a universal hazard pattern exists for clinical '
        'treatment dropout, and (4) to derive implications for the design of phase-specific '
        'retention interventions.')
    set_paragraph_spacing(p, space_after=12)
    
    # ---- METHODS ----
    doc.add_heading('2. Methods', level=1)
    
    doc.add_heading('2.1 Study Design and Data Sources', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'This is a secondary analysis of published retention/survival data from peer-reviewed '
        'literature. We identified major cohort studies, randomised controlled trials, and '
        'meta-analyses reporting Kaplan-Meier survival curves or tabulated retention data for '
        'each of the five clinical domains. Data points were extracted from published figures '
        'using digital curve extraction or from tabulated results reported in the text.')
    set_paragraph_spacing(p, space_after=6)
    
    # Data sources by domain
    doc.add_heading('2.1.1 HIV/ART Treatment Discontinuation', level=3)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Retention on ART data were obtained from the International Epidemiology Databases to '
        'Evaluate AIDS (IeDEA) consortium, which maintains observational cohorts across seven '
        'global regions.{10,11} We extracted regional retention curves for Sub-Saharan Africa '
        '(N = 256,000; Brinkhof 2009; Fox & Rosen 2010),{12,13} Asia-Pacific (N = 48,000; '
        'Zhou 2012),{14} Latin America (N = 35,000; CCASAnet, Wolff 2008),{15} North America '
        '(N = 82,000; NA-ACCORD),{16} and Europe (N = 95,000; COHERE).{17} Total combined '
        'sample: 516,000 patients. Follow-up ranged from 1 to 60 months from ART initiation.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('2.1.2 Antipsychotic Medication Discontinuation', level=3)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Data were drawn from three landmark studies: (1) the Clinical Antipsychotic Trials of '
        'Intervention Effectiveness (CATIE) Phase 1, a NIMH-funded pragmatic trial of 1,493 '
        'patients with chronic schizophrenia randomised to olanzapine, quetiapine, risperidone, '
        'perphenazine, or ziprasidone (Lieberman et al. 2005);{18} (2) the European First-Episode '
        'Schizophrenia Trial (EUFEST; N = 498; Kahn et al. 2008);{19} and (3) a Finnish '
        'nationwide register-based cohort study of 62,250 patients comparing long-acting injectable '
        '(LAI) versus oral antipsychotics (Tiihonen et al. 2017).{20} Total combined sample: '
        '63,682 patients.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('2.1.3 Substance Use Disorder Treatment Dropout', level=3)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Retention data were compiled from the Drug Abuse Treatment Outcome Studies (DATOS),{21} '
        'the Substance Abuse and Mental Health Services Administration Treatment Episode Data Set '
        '(SAMHSA TEDS),{22} Project MATCH/COMBINE for alcohol use disorder,{23} NIDA Clinical '
        'Trials Network cocaine studies,{24} and cannabis treatment trials.{25} Six substance-'
        'specific datasets were analysed: opioid use disorder (methadone maintenance, N = 12,500; '
        'buprenorphine, N = 8,200), alcohol (N = 15,000), cocaine (N = 6,800), cannabis '
        '(N = 4,500), and residential treatment for all substances (N = 28,000). Total: 75,000.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('2.1.4 Cardiac Rehabilitation Dropout', level=3)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Data were obtained from meta-analyses and large cohort studies: Turk-Adawi et al. '
        '(2013),{26} EUROASPIRE V (Kotseva et al. 2019),{27} Midence et al. (2020),{28} and '
        'Doll et al. (2015).{29} Five datasets were analysed by programme phase and indication: '
        'Phase I inpatient (N = 3,200), Phase II outpatient (N = 12,500), Phase III maintenance '
        '(N = 5,800), post-myocardial infarction (N = 8,500), and post-coronary artery bypass '
        'grafting (N = 4,200). Total: 34,200 patients.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('2.1.5 Clinical Trial Dropout', level=3)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Building on McChrystal et al. (2025), who demonstrated Weibull as the best-fitting '
        'distribution for clinical trial dropout across 90 RCTs (N = 86,107),{30} we analysed '
        'five therapeutic area subsets: oncology (N = 18,500), cardiovascular (N = 24,000), '
        'psychiatric (N = 15,200), diabetes/metabolic (N = 12,800), and vaccine trials '
        '(N = 28,500; including COVID-19 phase III).{31,32} Total: 99,000 participants.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('2.2 Weibull Distribution Fitting', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The two-parameter Weibull survival function was fitted to each dataset:\n\n'
        'S(t) = exp(-(t/\u03bb)^k)\n\n'
        'where k is the shape parameter and \u03bb is the scale parameter (characteristic life). '
        'Initial parameter estimates were obtained by linearising the Weibull survival function '
        'via the transformation ln(-ln(S(t))) = k\u00b7ln(t) - k\u00b7ln(\u03bb), followed by ordinary '
        'least-squares regression. These estimates were refined using Nelder-Mead nonlinear '
        'optimisation minimising the sum of squared errors between observed and predicted survival '
        'probabilities.{6,7}')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('2.3 Goodness-of-Fit Assessment', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Model fit was assessed using three metrics: (1) coefficient of determination (R\u00b2) '
        'comparing predicted versus observed survival probabilities, (2) Kolmogorov-Smirnov (KS) '
        'statistic measuring maximum absolute deviation between observed and predicted CDFs, and '
        '(3) visual assessment of Weibull probability plots (linearity of ln(-ln(S(t))) versus '
        'ln(t)). Bootstrap 95% confidence intervals for k and \u03bb were computed from 1,000 '
        'nonparametric resamples with added Gaussian noise (\u03c3 = 0.01) to account for curve '
        'digitisation uncertainty.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('2.4 Statistical Software', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'All analyses were performed in Python 3.11 using SciPy 1.11 for optimisation, '
        'NumPy 1.26 for numerical computation, and Matplotlib 3.8 for visualisation. '
        'Code is available at the project repository.')
    set_paragraph_spacing(p, space_after=12)
    
    # ---- RESULTS ----
    doc.add_heading('3. Results', level=1)
    
    doc.add_heading('3.1 Overall Findings', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'All 27 datasets across all five clinical domains yielded Weibull shape parameter '
        'k < 1, indicating a universal decreasing failure rate (DFR) pattern (Table 1). '
        'The overall range of k was 0.511 to 0.958, with a median of 0.689. The Weibull '
        'model demonstrated excellent fit across all datasets, with R\u00b2 values ranging from '
        '0.9807 to 0.9997 (median 0.9970) and KS statistics ranging from 0.0034 to 0.0392.')
    set_paragraph_spacing(p, space_after=6)
    
    # Table 1: Summary of Weibull parameters
    p = doc.add_paragraph()
    set_paragraph_spacing(p, space_before=12)
    run = p.add_run('Table 1. ')
    run.bold = True
    p.add_run('Summary of Weibull shape parameter (k) by domain.')
    
    table = doc.add_table(rows=6, cols=5)
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    headers = ['Domain', 'No. datasets', 'Total N', 'Mean k (range)', 'Median R\u00b2']
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True
    
    domain_data = [
        ('HIV/ART', 5, '516,000', '0.597 (0.582-0.627)', '0.9989'),
        ('Antipsychotic', 6, '63,682', '0.857 (0.683-0.958)', '0.9969'),
        ('Substance Use', 6, '75,000', '0.639 (0.511-0.733)', '0.9956'),
        ('Cardiac Rehab', 5, '34,200', '0.689 (0.613-0.775)', '0.9905'),
        ('Clinical Trial', 5, '99,000', '0.708 (0.679-0.742)', '0.9983'),
    ]
    for i, row_data in enumerate(domain_data):
        for j, val in enumerate(row_data):
            table.rows[i + 1].cells[j].text = str(val)
    
    set_paragraph_spacing(p, space_after=12)
    
    # Figure 1: Cross-domain k comparison
    p = doc.add_paragraph()
    set_paragraph_spacing(p, space_before=12)
    fig_path = os.path.join(FIG_DIR, 'fig11_cross_domain_k_comparison.png')
    if os.path.exists(fig_path):
        doc.add_picture(fig_path, width=Inches(6.5))
    
    p = doc.add_paragraph()
    set_paragraph_spacing(p, space_before=12)
    run = p.add_run('Fig. 1. ')
    run.bold = True
    p.add_run('Cross-domain comparison of Weibull shape parameters (k) with 95% bootstrap '
              'confidence intervals. The red dashed line indicates k = 1 (constant hazard). '
              'All 27 datasets fall in the k < 1 region, indicating decreasing failure rate '
              '(early dropout dominant).')
    set_paragraph_spacing(p, space_after=12)
    
    # Domain-specific results sections
    domain_sections = {
        'HIV/ART': {
            'heading': '3.2 HIV/ART Treatment Discontinuation',
            'fig_survival': 'fig1_hiv_art_treatment_discontinuation_survival.png',
            'fig_prob': 'fig2_hiv_art_treatment_discontinuation_probability.png',
            'fig_num_surv': 2, 'fig_num_prob': 3,
            'text': (
                'All five IeDEA regional cohorts showed consistently low k values ranging from '
                '0.582 (Latin America) to 0.627 (Sub-Saharan Africa), indicating a strong DFR '
                'pattern (Fig. 2). The remarkably narrow range across five continents (mean k = 0.597, '
                'SD = 0.018) suggests a highly conserved dropout hazard pattern in HIV/ART globally. '
                'The early treatment period (first 1-6 months) carries disproportionately high '
                'dropout risk across all regions. The scale parameter \u03bb showed greater regional '
                'variation (106.84 to 683.22 months), reflecting differences in absolute retention '
                'rates while maintaining similar temporal hazard shapes (Fig. 3).'
            ),
        },
        'Antipsychotic': {
            'heading': '3.3 Antipsychotic Medication Discontinuation',
            'fig_survival': 'fig3_antipsychotic_medication_discontinuation_survival.png',
            'fig_prob': 'fig4_antipsychotic_medication_discontinuation_probability.png',
            'fig_num_surv': 4, 'fig_num_prob': 5,
            'text': (
                'Antipsychotic discontinuation showed the highest k values among all five domains '
                '(mean = 0.857, range 0.683-0.958), approaching but not exceeding k = 1 (Fig. 4). '
                'The CATIE trial drugs showed near-constant hazard (k = 0.916-0.958), while '
                'the Finnish nationwide cohort showed lower k for both LAI (0.737) and oral (0.683) '
                'formulations (Fig. 5). The distinction between trial settings (CATIE, EUFEST: '
                'k ~ 0.92) and real-world cohorts (Finland: k ~ 0.71) suggests that clinical '
                'trial monitoring partially attenuates the early-dropout pattern. LAI formulations '
                'showed higher k than oral (0.737 vs 0.683), consistent with their design purpose '
                'of reducing early discontinuation.'
            ),
        },
        'Substance Use': {
            'heading': '3.4 Substance Use Disorder Treatment Dropout',
            'fig_survival': 'fig5_substance_use_disorder_treatment_dropout_survival.png',
            'fig_prob': 'fig6_substance_use_disorder_treatment_dropout_probability.png',
            'fig_num_surv': 6, 'fig_num_prob': 7,
            'text': (
                'Substance use disorder treatment showed a wide range of k values (0.511-0.733), '
                'with the steepest early dropout in residential programmes (k = 0.511) and '
                'opioid methadone maintenance (k = 0.576) (Fig. 6). Cannabis (k = 0.733) and '
                'alcohol outpatient (k = 0.714) showed relatively higher k values, indicating a '
                'less extreme DFR pattern (Fig. 7). The substance hierarchy of k values suggests '
                'that treatment modalities with higher external structure (residential) or '
                'pharmacological dependence management (methadone) paradoxically show steeper '
                'early attrition, potentially reflecting population-level severity differences '
                'rather than treatment effects per se.'
            ),
        },
        'Cardiac Rehab': {
            'heading': '3.5 Cardiac Rehabilitation Dropout',
            'fig_survival': 'fig7_cardiac_rehabilitation_dropout_survival.png',
            'fig_prob': 'fig8_cardiac_rehabilitation_dropout_probability.png',
            'fig_num_surv': 8, 'fig_num_prob': 9,
            'text': (
                'Cardiac rehabilitation dropout showed k values ranging from 0.613 (Phase II '
                'outpatient) to 0.775 (Phase I inpatient) (Fig. 8). The Phase I inpatient '
                'programme, being the most structured and closely monitored, showed the highest '
                'k (least steep early dropout), while Phase II outpatient showed the lowest, '
                'consistent with the transition from supervised to self-directed participation '
                'being a critical vulnerability period (Fig. 9). Post-CABG patients (k = 0.722) '
                'showed slightly higher k than post-MI patients (k = 0.664), possibly reflecting '
                'the more definitive nature of surgical intervention on patient engagement.'
            ),
        },
        'Clinical Trial': {
            'heading': '3.6 Clinical Trial Dropout',
            'fig_survival': 'fig9_clinical_trial_participant_dropout_survival.png',
            'fig_prob': 'fig10_clinical_trial_participant_dropout_probability.png',
            'fig_num_surv': 10, 'fig_num_prob': 11,
            'text': (
                'Clinical trial dropout showed relatively consistent k values across therapeutic '
                'areas (mean = 0.708, range 0.679-0.742), with the exception of the wider CI '
                'for vaccine trials (Fig. 10). Oncology trials showed the highest k (0.742), '
                'while cardiovascular trials showed the lowest (0.689) despite having the '
                'highest retention rates overall (Fig. 11). The narrow k range across diverse '
                'therapeutic areas suggests that the trial participation context itself (informed '
                'consent, monitoring, follow-up visits) imposes a relatively uniform temporal '
                'dropout pattern, consistent with the findings of McChrystal et al. (2025).{30}'
            ),
        },
    }
    
    fig_counter = 2  # Fig 1 already used
    for domain_key, info in domain_sections.items():
        doc.add_heading(info['heading'], level=2)
        p = doc.add_paragraph()
        add_text_with_citations(p, info['text'])
        set_paragraph_spacing(p, space_after=6)
        
        # Survival figure
        fig_path = os.path.join(FIG_DIR, info['fig_survival'])
        if os.path.exists(fig_path):
            doc.add_picture(fig_path, width=Inches(6.5))
        
        p = doc.add_paragraph()
        set_paragraph_spacing(p, space_before=12)
        run = p.add_run(f'Fig. {info["fig_num_surv"]}. ')
        run.bold = True
        p.add_run(f'{domain_key}: Retention curves with Weibull fits (left) and '
                  f'hazard functions (right).')
        set_paragraph_spacing(p, space_after=6)
        
        # Probability plot
        fig_path = os.path.join(FIG_DIR, info['fig_prob'])
        if os.path.exists(fig_path):
            doc.add_picture(fig_path, width=Inches(5.5))
        
        p = doc.add_paragraph()
        set_paragraph_spacing(p, space_before=12)
        run = p.add_run(f'Fig. {info["fig_num_prob"]}. ')
        run.bold = True
        p.add_run(f'{domain_key}: Weibull probability plot showing linearised goodness-of-fit.')
        set_paragraph_spacing(p, space_after=12)
    
    # Hazard taxonomy figure
    fig_path = os.path.join(FIG_DIR, 'fig12_hazard_taxonomy.png')
    if os.path.exists(fig_path):
        doc.add_picture(fig_path, width=Inches(6.5))
    
    p = doc.add_paragraph()
    set_paragraph_spacing(p, space_before=12)
    run = p.add_run('Fig. 12. ')
    run.bold = True
    p.add_run('Hazard pattern taxonomy: domain-specific hazard functions and summary comparison. '
              'All domains show decreasing hazard rates, with antipsychotic discontinuation '
              'nearest to constant hazard.')
    set_paragraph_spacing(p, space_after=12)
    
    # ---- DISCUSSION ----
    doc.add_heading('4. Discussion', level=1)
    
    doc.add_heading('4.1 Universal Decreasing Failure Rate Pattern', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The most striking finding of this study is the universality of the DFR pattern '
        '(k < 1) across all five clinical domains, all 27 datasets, and a combined sample '
        'exceeding 787,000 patients and trial participants. This consistency is remarkable '
        'given the diversity of conditions, treatments, populations, healthcare settings, and '
        'geographic regions represented. The DFR pattern indicates that dropout risk is highest '
        'in the early treatment period and decreases monotonically over time \u2014 the clinical '
        'analogue of "infant mortality" in reliability engineering.{7}')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'This finding has important theoretical implications. It suggests that treatment dropout '
        'is not primarily driven by cumulative treatment burden ("wear-out"), which would produce '
        'k > 1, nor by random events (k = 1), but rather by a selection process in which '
        'patients with the highest intrinsic dropout propensity leave earliest.{33} Those who '
        'survive the initial high-risk period become progressively "hardened" \u2014 a phenomenon '
        'well-characterised in reliability theory as the "screening" or "burn-in" effect.{7}')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('4.2 Domain-Specific Hazard Signatures', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Despite the universal DFR pattern, the magnitude of k varies meaningfully across domains. '
        'Antipsychotic medication showed the highest k values (mean 0.857), approaching constant '
        'hazard. This may reflect the chronic, relapsing nature of schizophrenia and the gradual '
        'emergence of side effects (e.g., metabolic syndrome, akathisia) that maintain ongoing '
        'dropout risk beyond the initial period.{18,20} In contrast, HIV/ART and substance use '
        'disorder showed the lowest k values (means 0.597 and 0.639, respectively), indicating '
        'steep early dropout gradients. For HIV/ART, this aligns with well-documented barriers '
        'in the first months of treatment including immune reconstitution inflammatory syndrome, '
        'pill burden adjustment, and disclosure-related psychosocial stressors.{12,13}')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('4.3 Clinical Implications: Front-Loading Retention Interventions', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'The universal DFR finding provides a quantitative rationale for front-loading retention '
        'interventions. Current practice often distributes retention resources uniformly over the '
        'treatment period or intensifies support only after patients miss appointments.{34} '
        'The Weibull framework suggests that the optimal strategy is to concentrate resources '
        'during the early treatment phase when hazard is highest. Specifically, for domains with '
        'the lowest k values (HIV/ART: k \u2248 0.60, substance use: k \u2248 0.64), the first 1-3 '
        'months represent the critical window where intensive support (frequent contact, peer '
        'navigation, digital adherence monitoring) would have the greatest impact per unit '
        'of investment.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('4.4 Comparison with Previous Weibull Applications', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Our previous application of Weibull analysis to tuberculosis treatment dropout '
        'yielded k = 1.22-1.31,{9} which is an IFR (increasing failure rate) pattern \u2014 '
        'in notable contrast to the universal DFR found in the present study. This difference '
        'is interpretable: TB treatment has a fixed 6-month duration with known endpoints, '
        'and dropout may indeed reflect cumulative treatment fatigue. In contrast, the five '
        'domains studied here predominantly involve open-ended or long-term treatments without '
        'a defined endpoint, where the DFR pattern emerges naturally from heterogeneous population '
        'selection. This distinction between fixed-duration treatments (IFR-prone) and open-ended '
        'treatments (DFR-prone) represents a potentially generalisable principle worthy of '
        'further investigation.')
    set_paragraph_spacing(p, space_after=6)
    
    doc.add_heading('4.5 Limitations', level=2)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'Several limitations merit discussion. First, data were reconstructed from published '
        'survival curves rather than individual patient data, which limits precision and '
        'precludes covariate-adjusted analysis. Second, the two-parameter Weibull model '
        'assumes a monotonic hazard function; mixed or non-monotonic patterns (e.g., bathtub '
        'curves) would require more complex models such as the three-parameter or mixture '
        'Weibull. Third, definitions of dropout/LTFU vary across studies and domains, '
        'introducing measurement heterogeneity. Fourth, publication bias toward studies with '
        'adequate follow-up may underestimate early dropout in some settings.')
    set_paragraph_spacing(p, space_after=12)
    
    # ---- CONCLUSIONS ----
    doc.add_heading('5. Conclusions', level=1)
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'This five-domain comparative study demonstrates that clinical treatment dropout follows '
        'a universal decreasing failure rate (DFR) pattern when modelled with the Weibull '
        'distribution. All 27 datasets spanning HIV/ART, antipsychotic medication, substance '
        'use disorder, cardiac rehabilitation, and clinical trials yielded shape parameter '
        'k < 1, with domain-specific signatures providing clinically interpretable hazard '
        'profiles. The reliability engineering framework offers a quantitative foundation for '
        'designing phase-specific retention interventions, with the strongest case for front-'
        'loading resources in domains with the steepest early dropout gradients (k \u2248 0.50-0.65). '
        'Future work should validate these findings with individual patient data and explore '
        'the fixed-duration versus open-ended treatment distinction as a predictor of hazard '
        'pattern directionality.')
    set_paragraph_spacing(p, space_after=12)
    
    # ---- REFERENCES ----
    doc.add_heading('References', level=1)
    references = [
        'World Health Organization. Adherence to long-term therapies: evidence for action. Geneva: WHO; 2003.',
        'Sabaté E. Adherence to long-term therapies: evidence for action. World Health Organization; 2003.',
        'Osterberg L, Blaschke T. Adherence to medication. N Engl J Med. 2005;353(5):487-497.',
        'DiMatteo MR. Variations in patients\' adherence to medical recommendations: a quantitative review of 50 years of research. Med Care. 2004;42(3):200-209.',
        'Cox DR. Regression models and life-tables. J R Stat Soc Series B. 1972;34(2):187-220.',
        'Weibull W. A statistical distribution function of wide applicability. J Appl Mech. 1951;18(3):293-297.',
        'Abernethy RB. The New Weibull Handbook. 5th ed. North Palm Beach, FL: Robert B. Abernethy; 2006.',
        'Onishi T. Weibull reliability analysis of sovereign reign durations: a cross-cultural comparative study. [Manuscript submitted].',
        'Onishi T. Weibull reliability analysis of tuberculosis treatment dropout: characterising hazard patterns for phase-specific retention interventions. [Manuscript submitted].',
        'Egger M, Ekouevi DK, Smith C, et al. Cohort Profile: the international epidemiological databases to evaluate AIDS (IeDEA) in sub-Saharan Africa. Int J Epidemiol. 2012;41(5):1256-1264.',
        'Haas AD, Zaniber A, von Wyl V, et al. Retention in care during the first 3 years of antiretroviral therapy for women in Malawi\'s Option B+ programme. Lancet HIV. 2016;3(4):e175-e182.',
        'Brinkhof MW, Pujades-Rodriguez M, Egger M. Mortality of patients lost to follow-up in antiretroviral treatment programmes in resource-limited settings. PLoS One. 2009;4(6):e5790.',
        'Fox MP, Rosen S. Patient retention in antiretroviral therapy programs up to three years on treatment in sub-Saharan Africa, 2007-2009: systematic review. Trop Med Int Health. 2010;15(s1):1-15.',
        'Zhou J, Tanuma J, Chaiwarith R, et al. Loss to followup in HIV-infected patients from Asia-Pacific region: results from TAHOD. AIDS Res Treat. 2012;2012:375217.',
        'Wolff MJ, Beltran CJ, Vasquez P, et al. The effect of antiretroviral therapy on mortality among HIV-infected patients in Latin America and the Caribbean. AIDS. 2008;22(8):977-987.',
        'Althoff KN, Buchacz K, Hall HI, et al. U.S. trends in antiretroviral therapy use, HIV RNA plasma viral loads, and CD4 T-lymphocyte cell counts among HIV-infected persons, 2000 to 2008. Ann Intern Med. 2012;157(5):325-335.',
        'Collaboration of Observational HIV Epidemiological Research Europe (COHERE). Response to combination antiretroviral therapy: variation by age. AIDS. 2008;22(12):1463-1473.',
        'Lieberman JA, Stroup TS, McEvoy JP, et al. Effectiveness of antipsychotic drugs in patients with chronic schizophrenia. N Engl J Med. 2005;353(12):1209-1223.',
        'Kahn RS, Fleischhacker WW, Boter H, et al. Effectiveness of antipsychotic drugs in first-episode schizophrenia and schizophreniform disorder. Lancet. 2008;371(9618):1085-1097.',
        'Tiihonen J, Mittendorfer-Rutz E, Majak M, et al. Real-world effectiveness of antipsychotic treatments in a nationwide cohort of 29,823 patients with schizophrenia. JAMA Psychiatry. 2017;74(7):686-693.',
        'Simpson DD, Joe GW, Brown BS. Treatment retention and follow-up outcomes in the Drug Abuse Treatment Outcome Study (DATOS). Psychol Addict Behav. 1997;11(4):294-307.',
        'Substance Abuse and Mental Health Services Administration. Treatment Episode Data Set (TEDS): 2017-2022. Rockville, MD: SAMHSA; 2023.',
        'Anton RF, O\'Malley SS, Ciraulo DA, et al. Combined pharmacotherapies and behavioral interventions for alcohol dependence: the COMBINE study. JAMA. 2006;295(17):2003-2017.',
        'National Institute on Drug Abuse. CTN dissemination library. National Institutes of Health; 2023.',
        'Budney AJ, Moore BA, Rocha HL, Higgins ST. Clinical trial of abstinence-based vouchers and cognitive-behavioral therapy for cannabis dependence. J Consult Clin Psychol. 2006;74(2):307-316.',
        'Turk-Adawi K, Sarrafzadegan N, Grace SL. Global availability of cardiac rehabilitation. Nat Rev Cardiol. 2014;11(10):586-596.',
        'Kotseva K, De Backer G, De Bacquer D, et al. Lifestyle and impact on cardiovascular risk factor control in coronary patients across 27 countries: results from the European Society of Cardiology ESC-EORP EUROASPIRE V registry. Eur J Prev Cardiol. 2019;26(8):824-835.',
        'Midence L, Arthur HM, Oh P, et al. Women\'s health behaviours and psychosocial well-being by cardiac rehabilitation program model. Can J Cardiol. 2016;32(8):S267-S268.',
        'Doll JA, Hsu JC, Tyan DB, et al. Cardiac rehabilitation completion and referral among women and men. J Am Coll Cardiol. 2015;65(10S):A1539.',
        'McChrystal R, Cro S, Carpenter JR. Best-fitting distributions for dropout in randomised clinical trials. BMC Med Res Methodol. 2025;25:42.',
        'Polack FP, Thomas SJ, Kitchin N, et al. Safety and efficacy of the BNT162b2 mRNA Covid-19 vaccine. N Engl J Med. 2020;383(27):2603-2615.',
        'Baden LR, El Sahly HM, Essink B, et al. Efficacy and safety of the mRNA-1273 SARS-CoV-2 vaccine. N Engl J Med. 2021;384(5):403-416.',
        'Vaupel JW, Manton KG, Stallard E. The impact of heterogeneity in individual frailty on the dynamics of mortality. Demography. 1979;16(3):439-454.',
        'Lawn SD, Myer L, Harling G, Orrell C, Bekker LG, Wood R. Determinants of mortality and nondeath losses from an antiretroviral treatment service in South Africa. BMJ. 2008;337:a1340.',
    ]
    
    for i, ref in enumerate(references, 1):
        p = doc.add_paragraph()
        add_text_with_citations(p, f'{i}. {ref}')
        set_paragraph_spacing(p, space_after=2)
    
    # Save
    output_path = os.path.join(MS_DIR, 'weibull_clinical_dropout_comprehensive_EN.docx')
    doc.save(output_path)
    print(f"  Saved: {output_path}")
    return output_path


# ============================================================
# Japanese manuscript
# ============================================================

def generate_comprehensive_manuscript_ja():
    """Generate comprehensive cross-domain manuscript in Japanese."""
    doc = Document()
    
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(11)
    
    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('臨床治療脱落のワイブル信頼性解析：\n'
                        '5領域横断比較研究による\n'
                        'フェーズ特異的維持介入のためのハザードパターン特性化')
    run.bold = True
    run.font.size = Pt(14)
    set_paragraph_spacing(title, space_after=12)
    
    authors = doc.add_paragraph()
    authors.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = authors.add_run('大西 辰樹')
    run.font.size = Pt(11)
    set_paragraph_spacing(authors, space_after=6)
    
    corr = doc.add_paragraph()
    corr.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = corr.add_run('責任著者: 大西辰樹 (bougtoir@gmail.com)')
    run.font.size = Pt(9)
    run.italic = True
    set_paragraph_spacing(corr, space_after=18)
    
    # Abstract
    h = doc.add_heading('抄録', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '背景: 治療脱落は臨床各領域における重大な課題であるが、脱落ハザードの時間的動態は'
        '十分に定量化されていない。信頼性工学のワイブル分布は、形状パラメータkを通じて'
        'ハザードパターンを特性化するパラメトリックフレームワークを提供する'
        '（k < 1: 減少型故障率[DFR]、k = 1: 一定ハザード、k > 1: 増加型故障率[IFR]）。')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '方法: 5つの臨床領域にわたる公表済みKaplan-Meier曲線から再構成した維持/生存データに'
        '2パラメータワイブル分布を適合させた。(1) HIV/ART治療中断（5地域、N=516,000）、'
        '(2) 統合失調症における抗精神病薬中断（6データセット、N=63,682）、'
        '(3) 物質使用障害治療脱落（6物質種、N=75,000）、'
        '(4) 心臓リハビリテーション脱落（5プログラムフェーズ、N=34,200）、'
        '(5) 臨床試験参加者脱落（5治療領域、N=99,000）。')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '結果: 全5領域・全27データセットでk < 1が得られ（範囲: 0.511-0.958）、普遍的なDFR'
        'パターンが示された。ワイブルモデルは優れた適合を示した（R\u00b2中央値=0.9970）。'
        '領域別平均k値: HIV/ART 0.597、抗精神病薬 0.857、物質使用 0.639、心臓リハ 0.689、'
        '臨床試験 0.708。')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '結論: 5つの異なる臨床領域にわたる普遍的DFRパターンは、脱落リスクが治療初期に最も高く'
        '時間とともに単調に減少することを示す。この知見は、初期治療フェーズにおける維持介入の'
        '前倒し集中配分を支持するものであり、信頼性工学フレームワークがフェーズ特異的'
        '維持戦略設計の定量的基盤を提供することを示している。')
    set_paragraph_spacing(p, space_after=12)
    
    p = doc.add_paragraph()
    run = p.add_run('キーワード: ')
    run.bold = True
    p.add_run('ワイブル分布; 信頼性工学; 治療脱落; 追跡脱落; ハザード関数; 生存時間解析; '
              'HIV/ART; 抗精神病薬; 物質使用障害; 心臓リハビリテーション; 臨床試験')
    set_paragraph_spacing(p, space_after=18)
    
    # Introduction
    doc.add_heading('1. 緒言', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '治療脱落は、追跡脱落（LTFU）、治療中断、不遵守とも呼ばれ、臨床医学において最も'
        '広範な課題の一つである。{1,2} 感染症管理から精神科治療まで、あらゆる治療領域において'
        '相当数の患者が治療目標達成前に治療から離脱する。従来の生存時間解析はリスク因子の'
        '同定に主眼を置いてきたが、脱落ハザードの時間的形状そのものの特性化は不十分であった。{5}')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '信頼性工学は、このギャップに対処する確立されたパラメトリックフレームワークを提供する。'
        '1951年にWeibullが開発したワイブル分布{6}は、形状パラメータkと尺度パラメータλで'
        '特性化される簡潔かつ柔軟なモデルである。kの値がハザードパターンを直接符号化する：'
        'k < 1はDFR（初期故障型）、k = 1は一定ハザード、k > 1はIFR（摩耗型）を示す。{7}')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '我々は先行研究において、このフレームワークを歴史的統治者の在位期間{8}および結核治療'
        '脱落{9}に適用した。本研究はこのアプローチを5つの主要臨床領域に拡張し、'
        '領域横断比較を行うことで、(1) 臨床治療脱落に普遍的ハザードパターンが存在するか、'
        '(2) 領域固有のハザード特性が同定できるか、(3) フェーズ特異的維持介入設計への'
        '示唆が得られるか、を検討する。')
    set_paragraph_spacing(p, space_after=12)
    
    # Methods
    doc.add_heading('2. 方法', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '本研究は公表済み文献からの維持/生存データの二次解析である。各領域の主要コホート研究、'
        'ランダム化比較試験、メタ解析からKaplan-Meier生存曲線または表形式維持データを抽出した。')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        'ワイブル生存関数 S(t) = exp(-(t/λ)^k) を各データセットに適合させた。初期推定は'
        '線形化変換 ln(-ln(S(t))) = k·ln(t) - k·ln(λ) による最小二乗回帰で取得し、'
        'Nelder-Mead法による非線形最適化で精緻化した。Bootstrap法（1,000回再抽出）で'
        '95%信頼区間を算出した。適合度はR\u00b2、KS統計量、ワイブル確率プロットで評価した。')
    set_paragraph_spacing(p, space_after=12)
    
    # Results
    doc.add_heading('3. 結果', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '全5領域・全27データセットにおいてk < 1が得られ（範囲: 0.511-0.958、中央値: 0.689）、'
        '普遍的なDFRパターンが確認された（表1）。R\u00b2中央値は0.9970であった。')
    set_paragraph_spacing(p, space_after=6)
    
    # Table
    p = doc.add_paragraph()
    set_paragraph_spacing(p, space_before=12)
    run = p.add_run('表1. ')
    run.bold = True
    p.add_run('領域別ワイブル形状パラメータ(k)の要約')
    
    table = doc.add_table(rows=6, cols=5)
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    headers_ja = ['領域', 'データセット数', '合計N', '平均k（範囲）', 'R\u00b2中央値']
    for i, h_text in enumerate(headers_ja):
        cell = table.rows[0].cells[i]
        cell.text = h_text
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True
    
    domain_data_ja = [
        ('HIV/ART', '5', '516,000', '0.597 (0.582-0.627)', '0.9989'),
        ('抗精神病薬', '6', '63,682', '0.857 (0.683-0.958)', '0.9969'),
        ('物質使用障害', '6', '75,000', '0.639 (0.511-0.733)', '0.9956'),
        ('心臓リハビリ', '5', '34,200', '0.689 (0.613-0.775)', '0.9905'),
        ('臨床試験', '5', '99,000', '0.708 (0.679-0.742)', '0.9983'),
    ]
    for i, row_data in enumerate(domain_data_ja):
        for j, val in enumerate(row_data):
            table.rows[i + 1].cells[j].text = str(val)
    
    set_paragraph_spacing(p, space_after=12)
    
    # Key figures (same images, Japanese captions)
    fig_path = os.path.join(FIG_DIR, 'fig11_cross_domain_k_comparison.png')
    if os.path.exists(fig_path):
        doc.add_picture(fig_path, width=Inches(6.5))
    
    p = doc.add_paragraph()
    set_paragraph_spacing(p, space_before=12)
    run = p.add_run('図1. ')
    run.bold = True
    p.add_run('ワイブル形状パラメータ(k)の領域横断比較。赤破線はk=1（一定ハザード）を示す。'
              '全27データセットがk < 1領域に位置し、DFR（初期脱落優勢）パターンを示す。')
    set_paragraph_spacing(p, space_after=12)
    
    fig_path = os.path.join(FIG_DIR, 'fig12_hazard_taxonomy.png')
    if os.path.exists(fig_path):
        doc.add_picture(fig_path, width=Inches(6.5))
    
    p = doc.add_paragraph()
    set_paragraph_spacing(p, space_before=12)
    run = p.add_run('図2. ')
    run.bold = True
    p.add_run('ハザードパターン分類学：領域別ハザード関数と要約比較。'
              '全領域で減少型ハザードを示し、抗精神病薬中断が一定ハザードに最も近い。')
    set_paragraph_spacing(p, space_after=12)
    
    # Discussion
    doc.add_heading('4. 考察', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '本研究の最も顕著な発見は、全5臨床領域・全27データセット・合計787,000人以上の患者・'
        '試験参加者にわたるDFRパターンの普遍性である。この一貫性は、疾患、治療、集団、'
        '医療環境、地理的地域の多様性を考慮すると注目に値する。DFRパターンは、脱落リスクが'
        '治療初期に最も高く時間とともに単調に減少することを示す――信頼性工学における'
        '「初期故障」の臨床的対応物である。{7}')
    set_paragraph_spacing(p, space_after=6)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '注目すべきは、先行する結核治療脱落の解析ではk=1.22-1.31のIFRパターンが得られた'
        'ことである。{9} この対比は解釈可能である：結核治療は6か月の固定期間であり、'
        '累積的治療負担（「摩耗」）が脱落を駆動しうる。一方、本研究の5領域は主に'
        'オープンエンドまたは長期治療であり、母集団の不均一性からの選択効果によって'
        'DFRパターンが自然に出現する。この固定期間治療（IFR傾向）vs オープンエンド治療'
        '（DFR傾向）の区別は、一般化可能な原理を示唆している。')
    set_paragraph_spacing(p, space_after=12)
    
    # Conclusion
    doc.add_heading('5. 結論', level=1)
    
    p = doc.add_paragraph()
    add_text_with_citations(p,
        '本5領域横断比較研究は、臨床治療脱落がワイブル分布による普遍的DFRパターンに'
        '従うことを実証した。全27データセットでk < 1が得られ、脱落リスクは治療初期に最も高い。'
        'この知見は、初期治療フェーズへの維持介入資源の前倒し集中配分を支持する。'
        '信頼性工学フレームワークは、フェーズ特異的維持戦略設計の定量的基盤を提供する。')
    set_paragraph_spacing(p, space_after=12)
    
    output_path = os.path.join(MS_DIR, 'weibull_clinical_dropout_comprehensive_JA.docx')
    doc.save(output_path)
    print(f"  Saved: {output_path}")
    return output_path


# ============================================================
# PPTX figure slides
# ============================================================

def generate_figure_pptx():
    """Generate editable PowerPoint with all figures (1 per slide, widescreen)."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    figures = [
        ('fig11_cross_domain_k_comparison.png', 
         'Fig. 1: Cross-Domain k Comparison',
         'Weibull shape parameters (k) with 95% CI across all 27 datasets and 5 domains.'),
        ('fig12_hazard_taxonomy.png',
         'Fig. 12: Hazard Pattern Taxonomy',
         'Domain-specific hazard functions and summary comparison.'),
        ('fig1_hiv_art_treatment_discontinuation_survival.png',
         'Fig. 2: HIV/ART Retention & Hazard',
         'Retention curves with Weibull fits and hazard functions, 5 IeDEA regions.'),
        ('fig2_hiv_art_treatment_discontinuation_probability.png',
         'Fig. 3: HIV/ART Weibull Probability Plot',
         'Linearised goodness-of-fit assessment for HIV/ART datasets.'),
        ('fig3_antipsychotic_medication_discontinuation_survival.png',
         'Fig. 4: Antipsychotic Retention & Hazard',
         'CATIE, EUFEST, and Finnish cohort retention curves and hazard functions.'),
        ('fig4_antipsychotic_medication_discontinuation_probability.png',
         'Fig. 5: Antipsychotic Weibull Probability Plot',
         'Linearised goodness-of-fit for antipsychotic datasets.'),
        ('fig5_substance_use_disorder_treatment_dropout_survival.png',
         'Fig. 6: Substance Use Retention & Hazard',
         'Substance-specific retention curves and hazard functions.'),
        ('fig6_substance_use_disorder_treatment_dropout_probability.png',
         'Fig. 7: Substance Use Weibull Probability Plot',
         'Linearised goodness-of-fit for substance use datasets.'),
        ('fig7_cardiac_rehabilitation_dropout_survival.png',
         'Fig. 8: Cardiac Rehab Retention & Hazard',
         'Phase-specific and indication-specific retention and hazard.'),
        ('fig8_cardiac_rehabilitation_dropout_probability.png',
         'Fig. 9: Cardiac Rehab Weibull Probability Plot',
         'Linearised goodness-of-fit for cardiac rehab datasets.'),
        ('fig9_clinical_trial_participant_dropout_survival.png',
         'Fig. 10: Clinical Trial Retention & Hazard',
         'Therapeutic area-specific trial dropout and hazard functions.'),
        ('fig10_clinical_trial_participant_dropout_probability.png',
         'Fig. 11: Clinical Trial Weibull Probability Plot',
         'Linearised goodness-of-fit for clinical trial datasets.'),
    ]
    
    for fig_file, title_text, caption_text in figures:
        fig_path = os.path.join(FIG_DIR, fig_file)
        if not os.path.exists(fig_path):
            continue
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2), 
                                          PptxInches(12.333), PptxInches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PptxPt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Image (centered)
        img_width = PptxInches(11)
        img_height = PptxInches(5.5)
        left = (prs.slide_width - img_width) // 2
        top = PptxInches(1.0)
        slide.shapes.add_picture(fig_path, left, top, img_width, img_height)
        
        # Caption
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.7),
                                          PptxInches(12.333), PptxInches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = caption_text
        p.font.size = PptxPt(12)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    output_path = os.path.join(MS_DIR, 'weibull_clinical_dropout_figures_EN.pptx')
    prs.save(output_path)
    print(f"  Saved: {output_path}")
    return output_path


# ============================================================
# Main
# ============================================================

def main():
    print("Generating manuscripts...")
    print("\n1. English comprehensive manuscript:")
    generate_comprehensive_manuscript_en()
    
    print("\n2. Japanese comprehensive manuscript:")
    generate_comprehensive_manuscript_ja()
    
    print("\n3. English figure slides (PPTX):")
    generate_figure_pptx()
    
    print("\nAll manuscripts generated successfully.")


if __name__ == '__main__':
    main()
